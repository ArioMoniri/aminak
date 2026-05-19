#!/usr/bin/env python3
"""Phase 14 — v2 redesign of the PPTX presentation.

Goals (per user feedback):
  - Less AI-generated look (no accent lines, asymmetric layouts, more whitespace)
  - Inter / Inter SemiBold font (with system sans-serif fallback)
  - More plots than text
  - Fewer words, tighter bullets
  - Cover ALL phases briefly, not only Phase 14

Outputs: /Users/ario/Downloads/Aminak Phase 14 Summary.pptx  (overwrites the user's copy)
         14_inhibitor_design/presentation/aminak_phase14_summary.pptx (repo copy)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import shutil

# Palette — restrained, scientific: charcoal + warm beige + single coral accent
INK     = RGBColor(0x14, 0x18, 0x1F)   # near-black titles
DARK    = RGBColor(0x21, 0x29, 0x33)   # dark bg slides
CHAR    = RGBColor(0x3A, 0x42, 0x4E)   # body text
MUTE    = RGBColor(0x8A, 0x94, 0xA1)   # captions
BONE    = RGBColor(0xF7, 0xF4, 0xEE)   # warm off-white bg
PAPER   = RGBColor(0xFD, 0xFB, 0xF7)   # paper white
RULE    = RGBColor(0xD8, 0xD2, 0xC6)   # thin rule lines
ACCENT  = RGBColor(0xE6, 0x55, 0x3F)   # coral — the single accent
GREEN   = RGBColor(0x3E, 0x7B, 0x4F)   # for PASS verdicts
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Inter SemiBold"  # falls back to system sans-serif
BODY = "Inter"
MONO = "JetBrains Mono"  # falls back to Menlo/Consolas

REPO = Path(__file__).resolve().parents[2]
OUT_REPO = REPO / "14_inhibitor_design" / "presentation" / "aminak_phase14_summary.pptx"
OUT_USER = Path("/Users/ario/Downloads/Aminak Phase 14 Summary.pptx")

# Image inventory
IMG = {
    "workflow":     REPO / "workflow_diagram_v3.png",
    "conservation_phase1": REPO / "01b_msa_v2" / "conservation_plot.png",
    "overlap":      REPO / "02_active_site" / "overlap_figure.png",
    "pymol_overview": REPO / "11_enhanced" / "pymol" / "wt_apo_overview.png",
    "pymol_close":  REPO / "04_pymol" / "02_closeup.png",
    "pymol_cavity": REPO / "04_pymol" / "04_cavity.png",
    "pymol_holo":   REPO / "11_enhanced" / "pymol" / "wt_holo_overview.png",
    "delta_apo":    REPO / "08d_analysis_v4" / "delta_vina_apo_holo.png",
    "delta_by_cat": REPO / "08d_analysis_v4" / "delta_vina_by_category.png",
    "phase8_flex":  REPO / "13_phase8" / "02_flexres" / "flex_vs_rigid.png",
    "phase8_alt":   REPO / "13_phase8" / "01_alt_scoring" / "alt_scoring_compare.png",
    "rama_best":    REPO / "10b_modeller_refined" / "04_refined_lovell" / "ramachandran_lovell_best_refined.png",
    "phase14_dist": REPO / "14_inhibitor_design" / "figures" / "fig1_distributions.png",
    "phase14_delta":REPO / "14_inhibitor_design" / "figures" / "fig2_delta_ranking.png",
    "phase14_apoholo": REPO / "14_inhibitor_design" / "figures" / "fig3_apo_holo_gap.png",
    "phase14_tier": REPO / "14_inhibitor_design" / "figures" / "fig4_tier_separation.png",
    "pose_indazole":  REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav18_CID7032.png",
    "pose_ibuprofen": REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav18_CID3672.png",
    "pose_indazole_cav2": REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav2_CID7032.png",
    "pose_flurbi":    REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav12_CID35814.png",
    "conservation_cav18": REPO / "14_inhibitor_design" / "04_allosteric" / "cavity18_evidence" / "figures" / "cavity18_conservation.png",
    "phylogeny_cav18":    REPO / "14_inhibitor_design" / "04_allosteric" / "cavity18_evidence" / "figures" / "cavity18_phylogeny_annot.png",
}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def text(slide, txt, x, y, w, h, size=14, bold=False, color=CHAR, align=PP_ALIGN.LEFT,
         font=BODY, italic=False, anchor=MSO_ANCHOR.TOP, char_spacing=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    paragraphs = txt if isinstance(txt, list) else [txt]
    for i, line in enumerate(paragraphs):
        if i > 0:
            p = tf.add_paragraph(); p.alignment = align
        if isinstance(line, list):
            for run_spec in line:
                r = p.add_run()
                r.text = run_spec["text"]
                r.font.size = Pt(run_spec.get("size", size))
                r.font.bold = run_spec.get("bold", bold)
                r.font.italic = run_spec.get("italic", italic)
                r.font.color.rgb = run_spec.get("color", color)
                r.font.name = run_spec.get("font", font)
        else:
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.bold = bold
            r.font.italic = italic; r.font.color.rgb = color
            r.font.name = font
    return tb

def add_image(slide, key, x, y, w=None, h=None):
    p = IMG[key]
    if not p.exists():
        print(f"  ! missing: {p}")
        return None
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    return slide.shapes.add_picture(str(p), x, y, **kw)

def thin_rule(slide, x, y, w, color=RULE, weight=0.75):
    """A thin horizontal rule — visual rhythm without being chrome-heavy."""
    sh = slide.shapes.add_connector(1, x, y, x + w, y)
    sh.line.color.rgb = color
    sh.line.width = Pt(weight)
    return sh

def slide_label(slide, label_text, color=ACCENT):
    """Tiny uppercase tracking label at top-left — magazine style."""
    text(slide, label_text.upper(),
         Inches(0.6), Inches(0.5), Inches(8), Inches(0.3),
         size=9, bold=True, color=color, font=HEAD)

def slide_number(slide, n, total):
    """Discreet page indicator."""
    text(slide, f"{n:02d} / {total:02d}",
         Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
         size=9, color=MUTE, font=MONO, align=PP_ALIGN.RIGHT)

TOTAL = 19

# ───────────────────────────────────────────────────────
# 01 — Title (dark, asymmetric)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, DARK)
# tiny project mark
text(slide, "AMINAK   /   PHASE 14",
     Inches(0.7), Inches(0.6), Inches(8), Inches(0.3),
     size=10, bold=True, color=RGBColor(0xC9,0xCC,0xD2), font=HEAD)
# huge title (asymmetric, left)
text(slide, "Designing inhibitors for TYMS\nat four binding sites.",
     Inches(0.7), Inches(2.4), Inches(11.5), Inches(2.6),
     size=50, bold=True, color=WHITE, font=HEAD)
# rule
thin_rule(slide, Inches(0.7), Inches(5.0), Inches(2.0), color=ACCENT, weight=1.5)
text(slide, "Active site · Cofactor site · Dimer interface · Allosteric",
     Inches(0.7), Inches(5.2), Inches(11.5), Inches(0.5),
     size=15, color=RGBColor(0xC9,0xCC,0xD2), font=BODY, italic=True)
# meta
text(slide, [
        [{"text": "github.com/ArioMoniri/aminak", "color": RGBColor(0xC9,0xCC,0xD2)}],
     ], Inches(0.7), Inches(6.85), Inches(10), Inches(0.3),
     size=10, font=MONO)
text(slide, "Ariorad Moniri   /   2026",
     Inches(0.7), Inches(7.15), Inches(10), Inches(0.3),
     size=9, color=MUTE, font=HEAD)

# ───────────────────────────────────────────────────────
# 02 — The question (3 stats, mostly whitespace)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, BONE)
slide_label(slide, "Setup")
slide_number(slide, 2, TOTAL)
text(slide, "The question.", Inches(0.7), Inches(1.0), Inches(12), Inches(1.0),
     size=44, bold=True, color=INK, font=HEAD)
text(slide, "What molecules will out-compete dUMP at the active site — or bind elsewhere on TYMS?",
     Inches(0.7), Inches(2.3), Inches(12), Inches(0.6),
     size=18, color=CHAR, font=BODY, italic=True)
thin_rule(slide, Inches(0.7), Inches(3.3), Inches(12), color=RULE)
# 3 large stats
def stat(slide, big, label, x, big_color=INK):
    text(slide, big, x, Inches(3.9), Inches(4.0), Inches(1.6),
         size=82, bold=True, color=big_color, font=HEAD)
    text(slide, label, x, Inches(5.7), Inches(4.0), Inches(0.8),
         size=13, color=CHAR, font=BODY)

stat(slide, "4",  "binding-site strategies", Inches(0.7))
stat(slide, "86", "docked compounds", Inches(4.7), big_color=INK)
stat(slide, "7",  "reviewer / corrector rounds", Inches(8.7), big_color=ACCENT)

text(slide, "TYMS · UniProt P04818 · PDB 1HVY · AutoDock Vina 1.2.7 · ±0.85 kcal/mol noise floor",
     Inches(0.7), Inches(6.9), Inches(12), Inches(0.3),
     size=10, color=MUTE, font=MONO)

# ───────────────────────────────────────────────────────
# 03 — Pipeline (visual: just the workflow diagram + short caption)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Pipeline · 8 phases")
slide_number(slide, 3, TOTAL)
text(slide, "From conservation to inhibitor.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=32, bold=True, color=INK, font=HEAD)
add_image(slide, "workflow", Inches(1.5), Inches(2.0), w=Inches(10.5))
text(slide, "Each phase reused by the next. Each independently reviewed by paired agents.",
     Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
     size=11, color=MUTE, font=BODY, italic=True, align=PP_ALIGN.CENTER)

# ───────────────────────────────────────────────────────
# 04 — Phase 1+2: Conservation
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 1+2 · MSA + active site")
slide_number(slide, 4, TOTAL)
text(slide, "The active site is highly conserved.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
add_image(slide, "conservation_phase1", Inches(0.7), Inches(2.0), w=Inches(8.5))
# Right column bullets
text(slide, ["10 verified TYMS orthologs",
             "JS-divergence per residue",
             "Cross-checked with UniProt + PDBe binding-site graph",
             "Catalytic Cys195, His196, R175/176/215 phosphate clamp"],
     Inches(9.5), Inches(2.3), Inches(3.5), Inches(4.5),
     size=14, color=CHAR, font=BODY)
text(slide, "→ defines the docking box used in every later phase.",
     Inches(9.5), Inches(5.5), Inches(3.5), Inches(0.6),
     size=12, italic=True, color=ACCENT, font=HEAD, bold=True)

# ───────────────────────────────────────────────────────
# 05 — Phase 3+4: Structure prep + renders
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 3+4 · Structure")
slide_number(slide, 5, TOTAL)
text(slide, "Dimer-aware receptor preparation.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
# 3-image strip
add_image(slide, "pymol_overview", Inches(0.7), Inches(2.0), w=Inches(4.0))
add_image(slide, "pymol_close",    Inches(4.85), Inches(2.0), w=Inches(4.0))
add_image(slide, "pymol_cavity",   Inches(9.0),  Inches(2.0), w=Inches(3.7))
# captions
text(slide, "Overview", Inches(0.7), Inches(5.5), Inches(4.0), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Active-site close-up", Inches(4.85), Inches(5.5), Inches(4.0), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Cavity carve-out", Inches(9.0), Inches(5.5), Inches(3.7), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, [
    "Both chains kept (active site spans the A↔B dimer).",
    "AMBER ff14SB charges via pdb2pqr30 (PROPKA pH 7.4).",
    "Cofactor (raltitrexed, D16) protonated to net −2 per copy.",
], Inches(0.7), Inches(6.05), Inches(12), Inches(1.2), size=12, color=CHAR, font=BODY,
   align=PP_ALIGN.CENTER)

# ───────────────────────────────────────────────────────
# 06 — Phase 5+6: WT docking
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 5+6 · WT Vina docking")
slide_number(slide, 6, TOTAL)
text(slide, "Wild-type docks at −8.78 kcal/mol.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
# big number left
text(slide, "−8.78", Inches(0.7), Inches(2.4), Inches(6.0), Inches(2.5),
     size=144, bold=True, color=INK, font=HEAD)
text(slide, "kcal/mol", Inches(0.7), Inches(5.0), Inches(6.0), Inches(0.5),
     size=20, color=MUTE, font=BODY, italic=True)
text(slide, "dUMP into 1HVY chain-A active site\nApo · exhaustiveness 32 · seed 42",
     Inches(0.7), Inches(5.7), Inches(6.0), Inches(1.4),
     size=14, color=CHAR, font=BODY)
# right: image of holo for visual rhythm
add_image(slide, "pymol_holo", Inches(7.5), Inches(2.0), w=Inches(5.3))
text(slide, "Holo state (cofactor present) →", Inches(7.5), Inches(6.4), Inches(5.3), Inches(0.4),
     size=12, italic=True, color=MUTE, font=BODY)

# ───────────────────────────────────────────────────────
# 07 — Phase 7: the null result
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 7 · mutagenesis", color=ACCENT)
slide_number(slide, 7, TOTAL)
text(slide, "Rigid Vina cannot resolve mutants.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
# plot left, big takeaway right
add_image(slide, "delta_apo", Inches(0.7), Inches(2.0), w=Inches(7.5))
# right column
text(slide, "+0.77", Inches(8.7), Inches(2.4), Inches(4.0), Inches(1.4),
     size=80, bold=True, color=ACCENT, font=HEAD)
text(slide, "largest holo Δ Vina (kcal/mol)",
     Inches(8.7), Inches(4.0), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "±0.85", Inches(8.7), Inches(4.7), Inches(4.0), Inches(1.0),
     size=60, bold=True, color=INK, font=HEAD)
text(slide, "Vina noise floor (Trott & Olson 2010)",
     Inches(8.7), Inches(6.0), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "→ honest null result.",
     Inches(8.7), Inches(6.6), Inches(4.0), Inches(0.4),
     size=14, italic=True, bold=True, color=ACCENT, font=HEAD)

# ───────────────────────────────────────────────────────
# 08 — Phase 8: smarter scoring
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 8 · Vinardo + flex-Vina")
slide_number(slide, 8, TOTAL)
text(slide, "Can smarter scoring fix it?",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
add_image(slide, "phase8_flex", Inches(0.7), Inches(2.0), w=Inches(6.0))
add_image(slide, "phase8_alt",  Inches(7.0), Inches(2.0), w=Inches(5.8))
# captions + 2-bullet takeaway
text(slide, "Flex-residue Vina", Inches(0.7), Inches(6.2), Inches(6.0), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Vinardo vs Vina", Inches(7.0), Inches(6.2), Inches(5.8), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Partial fix only. The R→E sign error needs proper PB electrostatics (MM-GBSA / FEP).",
     Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
     size=13, color=CHAR, font=BODY, italic=True, align=PP_ALIGN.CENTER)

# ───────────────────────────────────────────────────────
# 09 — Phase 6: Modeller homology
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 6 · Modeller homology")
slide_number(slide, 9, TOTAL)
text(slide, "Homology models beat the crystal on Ramachandran.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=26, bold=True, color=INK, font=HEAD)
add_image(slide, "rama_best", Inches(0.7), Inches(2.0), w=Inches(7.5))
# right column
text(slide, "95.4 %", Inches(8.7), Inches(2.4), Inches(4.0), Inches(1.4),
     size=72, bold=True, color=INK, font=HEAD)
text(slide, "Best refined model — favoured\n(Lovell 4-map partition)",
     Inches(8.7), Inches(4.0), Inches(4.0), Inches(0.9),
     size=13, color=CHAR, font=BODY)
text(slide, "92.2 %", Inches(8.7), Inches(5.1), Inches(4.0), Inches(1.0),
     size=48, bold=True, color=MUTE, font=HEAD)
text(slide, "1HVY crystal under same scheme",
     Inches(8.7), Inches(6.1), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)

# ───────────────────────────────────────────────────────
# 10 — Phase 14 section divider (dark)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, DARK)
text(slide, "PHASE 14", Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
     size=12, bold=True, color=ACCENT, font=HEAD)
text(slide, "Inhibitor design.",
     Inches(0.7), Inches(2.8), Inches(12), Inches(1.6),
     size=64, bold=True, color=WHITE, font=HEAD)
thin_rule(slide, Inches(0.7), Inches(4.5), Inches(3.0), color=ACCENT, weight=1.5)
text(slide, "Four orthogonal sites. 86 docked compounds. Two findings.",
     Inches(0.7), Inches(4.8), Inches(12), Inches(0.5),
     size=18, color=RGBColor(0xC9,0xCC,0xD2), italic=True, font=BODY)

# ───────────────────────────────────────────────────────
# 11 — Strategy 1: active-site (plot full)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 1 · active site")
slide_number(slide, 11, TOTAL)
text(slide, "Clean active-vs-prodrug gap. Kcal-noise silent.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=26, bold=True, color=INK, font=HEAD)
add_image(slide, "phase14_tier", Inches(0.7), Inches(2.0), w=Inches(8.0))
# right takeaways
text(slide, "−9.04", Inches(9.0), Inches(2.3), Inches(4.0), Inches(1.4),
     size=72, bold=True, color=INK, font=HEAD)
text(slide, "5-FdUMP   (canonical)",
     Inches(9.0), Inches(3.8), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "−4.95", Inches(9.0), Inches(4.6), Inches(4.0), Inches(1.4),
     size=72, bold=True, color=MUTE, font=HEAD)
text(slide, "5-FU   (prodrug)",
     Inches(9.0), Inches(6.1), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "≈3 kcal/mol enrichment signal.",
     Inches(9.0), Inches(6.6), Inches(4.0), Inches(0.4),
     size=11, italic=True, bold=True, color=ACCENT, font=HEAD)

# ───────────────────────────────────────────────────────
# 12 — Strategy 2: ★ Plevitrexed
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 2 · cofactor site", color=ACCENT)
slide_number(slide, 12, TOTAL)
text(slide, "Plevitrexed (ZD9331) above noise.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
# left: ranked-bar plot
add_image(slide, "phase14_delta", Inches(0.7), Inches(2.0), w=Inches(7.5))
# right: giant stat
text(slide, "−10.01", Inches(8.6), Inches(2.4), Inches(4.5), Inches(1.6),
     size=92, bold=True, color=ACCENT, font=HEAD)
text(slide, "kcal/mol", Inches(8.6), Inches(4.2), Inches(4.5), Inches(0.4),
     size=14, color=MUTE, italic=True, font=BODY)
thin_rule(slide, Inches(8.6), Inches(4.8), Inches(2.5), color=INK, weight=1.5)
text(slide, ["Δ −0.88 vs raltitrexed",
             "above the ±0.85 noise floor",
             "reproducible across both seeds",
             "consistent with Ki ~nM (Jackman 1997)"],
     Inches(8.6), Inches(5.0), Inches(4.5), Inches(2.0),
     size=13, color=CHAR, font=BODY)

# ───────────────────────────────────────────────────────
# 13 — Strategy 3: dimer interface (null)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 3 · dimer interface")
slide_number(slide, 13, TOTAL)
text(slide, "Documented null result.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
# Centre: pull-quote
quote_y = Inches(2.6)
text(slide, "Canonical LR-octapeptide   worse   than scrambled control.",
     Inches(0.7), quote_y, Inches(12), Inches(1.2),
     size=32, bold=True, color=INK, font=HEAD,
     align=PP_ALIGN.CENTER)
text(slide, "specificity vs scrambled = +1.48 kcal/mol",
     Inches(0.7), Inches(3.9), Inches(12), Inches(0.5),
     size=16, italic=True, color=ACCENT, font=BODY, align=PP_ALIGN.CENTER)
thin_rule(slide, Inches(5.0), Inches(4.7), Inches(3.3), color=RULE)
text(slide, ["Rigid Vina is the wrong engine for ≥6-mer peptides (Hassan 2017).",
             "HPEPDOCK web service unreachable at execution.",
             "Null is the correct conclusion — not a methodology failure."],
     Inches(2.5), Inches(5.1), Inches(8.5), Inches(2.2),
     size=14, color=CHAR, font=BODY)

# ───────────────────────────────────────────────────────
# 14 — Strategy 4: cavity 18 headline
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 4 · allosteric", color=ACCENT)
slide_number(slide, 14, TOTAL)
text(slide, "An under-explored druggable cavity.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
# 4 big stats in a row
def big_stat(x, big, label, color=INK):
    text(slide, big, x, Inches(2.4), Inches(3.0), Inches(1.6),
         size=70, bold=True, color=color, font=HEAD)
    text(slide, label, x, Inches(4.0), Inches(3.0), Inches(0.8),
         size=11, color=CHAR, font=BODY)

big_stat(Inches(0.7),  "0.994", "FPocket druggability\ncavity 18", color=ACCENT)
big_stat(Inches(3.8),  "0.828", "C2-symmetric mirror\ncavity 17", color=INK)
big_stat(Inches(6.9),  "−7.52", "1H-indazole\nkcal/mol", color=ACCENT)
big_stat(Inches(10.0), "−7.28", "ibuprofen\nkcal/mol", color=ACCENT)

thin_rule(slide, Inches(0.7), Inches(5.6), Inches(12), color=RULE)
text(slide, "Two unrelated drug-like fragments. Same pocket. Different chemistries.",
     Inches(0.7), Inches(5.9), Inches(12), Inches(0.5),
     size=18, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "≈ 2 kcal/mol better than anywhere else on the chain-A or chain-B surface.",
     Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
     size=13, color=MUTE, font=BODY, italic=True, align=PP_ALIGN.CENTER)

# ───────────────────────────────────────────────────────
# 15 — indazole pose (image dominant)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · top hit", color=ACCENT)
slide_number(slide, 15, TOTAL)
text(slide, "1H-indazole.",
     Inches(0.7), Inches(0.95), Inches(8), Inches(0.7),
     size=36, bold=True, color=INK, font=HEAD)
text(slide, "Kinase-inhibitor scaffold · −7.52 kcal/mol",
     Inches(0.7), Inches(1.7), Inches(8), Inches(0.5),
     size=14, color=MUTE, italic=True, font=BODY)
# image dominant left
add_image(slide, "pose_indazole", Inches(0.5), Inches(2.4), w=Inches(7.5))
# right: 5 key contacts as tight bullets
text(slide, "Engages", Inches(8.5), Inches(2.4), Inches(4.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, [
    [{"text": "Phe55  ", "bold": True, "color": INK}, {"text": "H-bond + π"}],
    [{"text": "Asn201  ", "bold": True, "color": INK}, {"text": "H-bond"}],
    [{"text": "Phe200  ", "bold": True, "color": INK}, {"text": "π-stack"}],
    [{"text": "Leu196   ", "bold": True, "color": ACCENT}, {"text": "allosteric loop 181-197 ★", "color": ACCENT}],
    [{"text": "Gly197   ", "bold": True, "color": ACCENT}, {"text": "allosteric loop 181-197 ★", "color": ACCENT}],
    [{"text": "Ile83  ", "bold": True, "color": INK}, {"text": "hydrophobic wall"}],
], Inches(8.5), Inches(2.85), Inches(4.5), Inches(3.5), size=14, color=CHAR, font=BODY)
thin_rule(slide, Inches(8.5), Inches(6.0), Inches(4.0), color=RULE)
text(slide, "3 of 10 contact residues are on the published allosteric communication loop (Anderson 2012, Pozzi 2019).",
     Inches(8.5), Inches(6.2), Inches(4.5), Inches(1.0),
     size=11, italic=True, color=ACCENT, font=BODY)

# ───────────────────────────────────────────────────────
# 16 — ibuprofen pose
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · second hit", color=ACCENT)
slide_number(slide, 16, TOTAL)
text(slide, "Ibuprofen.",
     Inches(0.7), Inches(0.95), Inches(8), Inches(0.7),
     size=36, bold=True, color=INK, font=HEAD)
text(slide, "NSAID, COX1/2 · −7.28 kcal/mol",
     Inches(0.7), Inches(1.7), Inches(8), Inches(0.5),
     size=14, color=MUTE, italic=True, font=BODY)
add_image(slide, "pose_ibuprofen", Inches(0.5), Inches(2.4), w=Inches(7.5))
text(slide, "Engages", Inches(8.5), Inches(2.4), Inches(4.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, [
    [{"text": "Lys283  ", "bold": True, "color": ACCENT}, {"text": "salt bridge ★", "color": ACCENT}],
    [{"text": "Lys52   ", "bold": True, "color": ACCENT}, {"text": "salt bridge ★", "color": ACCENT}],
    [{"text": "Phe200  ", "bold": True, "color": INK}, {"text": "π-stack"}],
    [{"text": "Leu196  ", "bold": True, "color": ACCENT}, {"text": "loop 181-197 ★", "color": ACCENT}],
    [{"text": "Phe55   ", "bold": True, "color": INK}, {"text": "hydrophobic"}],
    [{"text": "Ile83   ", "bold": True, "color": INK}, {"text": "hydrophobic"}],
], Inches(8.5), Inches(2.85), Inches(4.5), Inches(3.5), size=14, color=CHAR, font=BODY)
thin_rule(slide, Inches(8.5), Inches(6.0), Inches(4.0), color=RULE)
text(slide, "Double-lysine salt-bridge clamps the deprotonated carboxylate at pH 7.4.",
     Inches(8.5), Inches(6.2), Inches(4.5), Inches(0.8),
     size=11, italic=True, color=ACCENT, font=BODY)

# ───────────────────────────────────────────────────────
# 17 — Cavity 18 conservation
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · conservation", color=ACCENT)
slide_number(slide, 17, TOTAL)
text(slide, "Conserved at the contact face.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=28, bold=True, color=INK, font=HEAD)
add_image(slide, "conservation_cav18", Inches(0.7), Inches(2.0), w=Inches(8.5))
# right column: the 7
text(slide, "100% conserved", Inches(9.5), Inches(2.3), Inches(3.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, "across 10 orthologs",
     Inches(9.5), Inches(2.6), Inches(3.5), Inches(0.3),
     size=10, italic=True, color=MUTE, font=BODY)
text(slide, [
    [{"text": "Gly54", "bold": True, "color": INK}],
    [{"text": "Glu87", "bold": True, "color": INK}],
    [{"text": "Met190", "bold": True, "color": ACCENT}, {"text": "  loop", "color": MUTE, "size": 10}],
    [{"text": "Ala191", "bold": True, "color": ACCENT}, {"text": "  loop", "color": MUTE, "size": 10}],
    [{"text": "Leu196", "bold": True, "color": ACCENT}, {"text": "  loop", "color": MUTE, "size": 10}],
    [{"text": "Phe200", "bold": True, "color": INK}],
    [{"text": "Asn201", "bold": True, "color": INK}],
], Inches(9.5), Inches(3.1), Inches(3.5), Inches(2.8), size=14, color=CHAR, font=BODY)
thin_rule(slide, Inches(9.5), Inches(6.0), Inches(3.0), color=RULE)
text(slide, "6 of 7 are contact residues for at least one top hit.",
     Inches(9.5), Inches(6.2), Inches(3.5), Inches(1.0),
     size=11, italic=True, bold=True, color=ACCENT, font=BODY)

# ───────────────────────────────────────────────────────
# 18 — Cavity 18 phylogeny
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · phylogeny", color=ACCENT)
slide_number(slide, 18, TOTAL)
text(slide, "A possible species-selective handle.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=28, bold=True, color=INK, font=HEAD)
add_image(slide, "phylogeny_cav18", Inches(0.5), Inches(2.0), w=Inches(8.0))
# right column: mut counts
text(slide, "Mutations vs human", Inches(8.7), Inches(2.3), Inches(4.0), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
data = [
    ("Mus musculus", "2", GREEN),
    ("Rattus norvegicus", "2", GREEN),
    ("Drosophila, yeast, E. coli", "12", INK),
    ("L. casei", "16", INK),
    ("Bacteriophage T4", "18", INK),
    ("Plasmodium falciparum", "21", ACCENT),
]
y = 3.0
for sp, n, col in data:
    text(slide, sp, Inches(8.7), Inches(y), Inches(2.8), Inches(0.35),
         size=11, color=CHAR, font=BODY)
    text(slide, n, Inches(11.6), Inches(y-0.07), Inches(1.0), Inches(0.5),
         size=18 if col == ACCENT else 15, bold=True, color=col, font=HEAD, align=PP_ALIGN.RIGHT)
    y += 0.5
thin_rule(slide, Inches(8.7), Inches(6.4), Inches(4.0), color=RULE)
text(slide, "Plasmodium TYMS diverges enough at cavity 18 for plausible selectivity vs the conserved active site.",
     Inches(8.7), Inches(6.55), Inches(4.5), Inches(1.0),
     size=10, italic=True, color=ACCENT, font=BODY)

# ───────────────────────────────────────────────────────
# 19 — TL;DR (dark, asymmetric, single pull-quote)
# ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK); add_bg(slide, DARK)
text(slide, "TL;DR", Inches(0.7), Inches(0.6), Inches(8), Inches(0.4),
     size=12, bold=True, color=ACCENT, font=HEAD)
text(slide, "We didn't find a new clinical lead.",
     Inches(0.7), Inches(1.5), Inches(12), Inches(0.9),
     size=38, bold=True, color=WHITE, font=HEAD)
text(slide, "We found a plausible third TYMS binding site.",
     Inches(0.7), Inches(2.5), Inches(12), Inches(0.9),
     size=38, bold=True, color=ACCENT, font=HEAD)
thin_rule(slide, Inches(0.7), Inches(3.7), Inches(2.5), color=ACCENT, weight=1.5)
text(slide, [
    "FPocket druggability 0.994 + 0.828 C2 mirror.",
    "Indazole + ibuprofen dock at −7.5 kcal/mol.",
    "6 of 7 ultra-conserved residues are at the contact face.",
    "21-position divergence in Plasmodium → species-selective handle.",
    "Anionic head + aromatic body + hydrophobic tail = pharmacophore.",
], Inches(0.7), Inches(4.0), Inches(12), Inches(3.0), size=18, color=RGBColor(0xC9,0xCC,0xD2), font=BODY)
text(slide, "github.com/ArioMoniri/aminak",
     Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
     size=10, color=MUTE, font=MONO)

# ───────────────────────────────────────────────────────
# Save (both locations)
# ───────────────────────────────────────────────────────
OUT_REPO.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_REPO)
shutil.copyfile(OUT_REPO, OUT_USER)
print(f"→ {OUT_REPO}  ({OUT_REPO.stat().st_size/1024:.0f} KB)")
print(f"→ {OUT_USER}")
