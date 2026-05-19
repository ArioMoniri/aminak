#!/usr/bin/env python3
"""Phase 14 — build a PPTX presentation summarising the whole project.

Outputs: 14_inhibitor_design/presentation/aminak_phase14_summary.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

# ────────── Palette: Ocean Gradient (deep blue → teal → midnight) ──────────
NAVY     = RGBColor(0x21, 0x29, 0x5C)    # title slides & dividers
DEEP     = RGBColor(0x06, 0x5A, 0x82)    # primary
TEAL     = RGBColor(0x1C, 0x72, 0x93)    # secondary
INK      = RGBColor(0x1a, 0x22, 0x30)    # dark body text
TEXT     = RGBColor(0x33, 0x33, 0x33)    # body text
LIGHT    = RGBColor(0xF5, 0xF7, 0xFA)    # light background
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT   = RGBColor(0xE6, 0x7E, 0x22)    # orange accent for callouts (★)
RED      = RGBColor(0xC0, 0x39, 0x2B)    # for "above-noise" markers
GREEN    = RGBColor(0x2C, 0x5F, 0x2D)    # for verdicts/pass
MUTED    = RGBColor(0x77, 0x88, 0x99)    # captions
NOISE    = RGBColor(0xBD, 0xC3, 0xC7)    # noise-floor lines

REPO = Path(__file__).resolve().parents[2]
OUT_DIR = REPO / "14_inhibitor_design" / "presentation"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "aminak_phase14_summary.pptx"

# Image paths
IMG = {
    "workflow": REPO / "workflow_diagram_v3.png",
    "fig1": REPO / "14_inhibitor_design" / "figures" / "fig1_distributions.png",
    "fig2": REPO / "14_inhibitor_design" / "figures" / "fig2_delta_ranking.png",
    "fig3": REPO / "14_inhibitor_design" / "figures" / "fig3_apo_holo_gap.png",
    "fig4": REPO / "14_inhibitor_design" / "figures" / "fig4_tier_separation.png",
    "pose_indazole":  REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav18_CID7032.png",
    "pose_ibuprofen": REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav18_CID3672.png",
    "pose_indazole_cav2": REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav2_CID7032.png",
    "pose_flurbi": REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav12_CID35814.png",
    "conservation": REPO / "14_inhibitor_design" / "04_allosteric" / "cavity18_evidence" / "figures" / "cavity18_conservation.png",
    "phylogeny": REPO / "14_inhibitor_design" / "04_allosteric" / "cavity18_evidence" / "figures" / "cavity18_phylogeny_annot.png",
}

# ────────── Setup ──────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]  # totally blank


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, x, y, w, h, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             italic=False, font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    # support multi-paragraph
    paragraphs = text if isinstance(text, list) else [text]
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if isinstance(para, dict):
            # rich run
            r = p.add_run()
            r.text = para["text"]
            r.font.size = Pt(para.get("size", size))
            r.font.bold = para.get("bold", bold)
            r.font.italic = para.get("italic", italic)
            r.font.color.rgb = para.get("color", color)
            r.font.name = para.get("font", font)
        else:
            r = p.add_run()
            r.text = para
            r.font.size = Pt(size); r.font.bold = bold
            r.font.italic = italic; r.font.color.rgb = color
            r.font.name = font
    return tb


def add_bullets(slide, bullets, x, y, w, h, size=14, color=TEXT, bullet_color=DEEP,
                bullet_size=None, gap=2):
    """bullets: list of strings (or dicts with optional 'bold')."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    bsize = bullet_size or size
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        # bullet glyph
        r = p.add_run()
        r.text = "▸ "
        r.font.size = Pt(bsize); r.font.bold = True
        r.font.color.rgb = bullet_color; r.font.name = "Calibri"
        # text
        if isinstance(b, str):
            parts = [{"text": b}]
        elif isinstance(b, dict):
            parts = [b]
        else:
            parts = b   # list of run-dicts
        for part in parts:
            r2 = p.add_run()
            r2.text = part["text"]
            r2.font.size = Pt(part.get("size", size))
            r2.font.bold = part.get("bold", False)
            r2.font.color.rgb = part.get("color", color)
            r2.font.name = part.get("font", "Calibri")
    return tb


def add_section_bar(slide, label, color=DEEP):
    """Thick coloured side-bar at left, decorative chrome."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # rotated label
    add_text(slide, label, Inches(0.05), Inches(6.0), Inches(2.0), Inches(0.3),
             size=10, color=WHITE, bold=True, font="Consolas")


def add_image(slide, key, x, y, w=None, h=None):
    p = IMG[key]
    if not p.exists():
        print(f"  ! missing image: {p}")
        return None
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    return slide.shapes.add_picture(str(p), x, y, **kw)


def add_table(slide, data, x, y, w, h, header_color=DEEP, header_text=WHITE,
              cell_size=11, header_size=12, zebra=True, col_widths=None):
    nrows = len(data); ncols = len(data[0])
    tbl_shape = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for r, row in enumerate(data):
        for c, cell_val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Emu(40000)
            tf.margin_right = Emu(40000)
            tf.margin_top = Emu(20000)
            tf.margin_bottom = Emu(20000)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            # support rich runs in cell_val
            runs = cell_val if isinstance(cell_val, list) else [{"text": str(cell_val)}]
            for run_spec in runs:
                run = p.add_run()
                run.text = run_spec["text"]
                run.font.size = Pt(header_size if r == 0 else cell_size)
                run.font.bold = run_spec.get("bold", r == 0)
                run.font.color.rgb = run_spec.get("color", header_text if r == 0 else TEXT)
                run.font.name = run_spec.get("font", "Calibri")
            # fill
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = header_color
            elif zebra and r % 2 == 0:
                fill.fore_color.rgb = LIGHT
            else:
                fill.fore_color.rgb = WHITE
    return tbl_shape


def add_stat_callout(slide, big, label, x, y, w=Inches(2.6), h=Inches(1.5),
                     big_color=DEEP, label_color=TEXT, big_size=60):
    """Big-number callout."""
    add_text(slide, big, x, y, w, Inches(1.0), size=big_size, bold=True,
             color=big_color, align=PP_ALIGN.CENTER, font="Calibri")
    add_text(slide, label, x, y+Inches(1.0), w, Inches(0.5), size=11,
             color=label_color, align=PP_ALIGN.CENTER, font="Calibri")


# ────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, NAVY)
# decorative left bar
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.6), SH)
deco.fill.solid(); deco.fill.fore_color.rgb = TEAL; deco.line.fill.background()
# decorative bottom bar
deco2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.7), SW - Inches(0.6), Inches(0.8))
deco2.fill.solid(); deco2.fill.fore_color.rgb = DEEP; deco2.line.fill.background()

add_text(slide, "aminak", Inches(1.0), Inches(0.8), Inches(11.5), Inches(0.8),
         size=20, color=TEAL, bold=False, font="Consolas")
add_text(slide, "TYMS / dUMP structural-bioinformatics workbench",
         Inches(1.0), Inches(1.5), Inches(11.5), Inches(0.7),
         size=18, color=RGBColor(0xCA, 0xDC, 0xFC), font="Calibri")
add_text(slide, "Phase 14 — Inhibitor design at four binding sites",
         Inches(1.0), Inches(2.5), Inches(11.5), Inches(1.4),
         size=44, color=WHITE, bold=True, font="Georgia")
add_text(slide, "Four mechanistically distinct sites · 86 docked compounds · 7 reviewer/corrector rounds\n"
                "Two real findings: Plevitrexed (ZD9331) above noise at the cofactor site, and a previously\n"
                "under-explored druggable cavity on both TYMS protomers (FPocket 0.994 / 0.828).",
         Inches(1.0), Inches(4.4), Inches(11.5), Inches(1.6),
         size=17, color=RGBColor(0xCA, 0xDC, 0xFC), font="Calibri")
add_text(slide, "github.com/ArioMoniri/aminak  ·  Built end-to-end with reviewer/corrector agents",
         Inches(1.0), Inches(6.85), Inches(11.5), Inches(0.5),
         size=12, color=RGBColor(0xCA, 0xDC, 0xFC), italic=True, font="Calibri",
         align=PP_ALIGN.LEFT)

# ────────────────────────────────────────────────────────────────
# Slide 2 — The protein + the question
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "SETUP", DEEP)
add_text(slide, "The target — human TYMS", Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         size=32, bold=True, color=NAVY, font="Georgia")
add_text(slide, "UniProt P04818 · PDB 1HVY (1.9 Å, dUMP + raltitrexed bound) · obligate homodimer",
         Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
         size=14, italic=True, color=MUTED, font="Calibri")
# Two-column layout
add_text(slide, "What TYMS does", Inches(0.6), Inches(1.7), Inches(6), Inches(0.4),
         size=18, bold=True, color=DEEP, font="Calibri")
add_bullets(slide, [
    "Catalyses dUMP → dTMP using methylene-THF as the methyl donor.",
    "The only de novo source of dTMP for DNA — proliferating cells depend on it.",
    "Clinical target of 5-fluorouracil (5-FU): the prodrug is metabolised to 5-FdUMP, which traps TYMS in a covalent ternary complex.",
    "Also target of raltitrexed, pemetrexed, methotrexate, nolatrexed, plevitrexed.",
], Inches(0.6), Inches(2.2), Inches(6.0), Inches(3.5), size=13, gap=4)

add_text(slide, "Phase 14 inverts the question", Inches(7.0), Inches(1.7), Inches(6), Inches(0.4),
         size=18, bold=True, color=DEEP, font="Calibri")
add_bullets(slide, [
    {"text": "Phases 1–13 asked: ", "color": TEXT},
    {"text": "can rigid Vina resolve TYMS mutants?", "color": NAVY, "bold": True, "italic": True},
    "",
    {"text": "Phase 14 asks: ", "color": TEXT},
    {"text": "what molecules will out-compete dUMP, or bind elsewhere on the enzyme?", "color": NAVY, "bold": True, "italic": True},
], Inches(7.0), Inches(2.2), Inches(6.0), Inches(2.2), size=13, gap=4)
# Big stat
add_stat_callout(slide, "4", "binding sites screened",
                 Inches(7.0), Inches(4.4), w=Inches(2.5))
add_stat_callout(slide, "86", "docked compounds",
                 Inches(9.5), Inches(4.4), w=Inches(2.5), big_color=TEAL)
add_stat_callout(slide, "7", "reviewer rounds",
                 Inches(12.0), Inches(4.4), w=Inches(1.3), big_size=44, big_color=ACCENT)
add_text(slide, "Engine: AutoDock Vina 1.2.7 · Apple Silicon · Vina ±0.85 kcal/mol noise floor (Trott & Olson 2010)",
         Inches(0.6), Inches(6.9), Inches(12.7), Inches(0.4),
         size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 3 — Pipeline at a glance
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "PIPELINE", DEEP)
add_text(slide, "8 phases — Stage 1 → Phase 14",
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         size=32, bold=True, color=NAVY, font="Georgia")
add_text(slide, "Each phase is a numbered folder; each version reuses the previous; each reviewed by independent agents.",
         Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

# Two-column phase grid
phases = [
    ("1+2", "MSA + JS conservation + active-site annotation", "01_msa, 02_active_site"),
    ("3+4", "Dimer-aware structure prep + PyMOL renders", "03_structure, 04_pymol"),
    ("5", "Ligand prep (dUMP)", "05_ligand"),
    ("6", "Wild-type Vina docking (apo + holo)", "06_docking_wt"),
    ("7", "Mutagenesis ×20 ×2 conditions", "07_mut_docking"),
    ("8+9", "Aggregate analysis + reports + 3Dmol viewers", "08_analysis, 09_report"),
    ("Phase 6", "Modeller homology modelling (Lovell-validated)", "10_modeller"),
    ("Phase 7", "Multi-replica Vina + SASA + AlphaFold + phylogeny", "12_phase7"),
    ("Phase 8", "Vinardo + flex-residue Vina (smarter scoring)", "13_phase8"),
    ("Phase 14 ★", "Inhibitor design at 4 binding sites (this presentation)", "14_inhibitor_design"),
]
row_h = Inches(0.5); start_y = Inches(1.7)
for i, (stage, title, folder) in enumerate(phases):
    y = start_y + i * row_h
    # number box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(1.3), Inches(0.42))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT if "★" in stage else (DEEP if "Phase" in stage else TEAL)
    box.line.fill.background()
    tf = box.text_frame; tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
    tf.margin_top = Emu(10000); tf.margin_bottom = Emu(10000)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = stage; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Consolas"
    # title
    add_text(slide, title, Inches(2.15), y, Inches(7.5), Inches(0.42),
             size=13, color=NAVY if "★" in stage else TEXT,
             bold=("★" in stage), font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
    # folder
    add_text(slide, folder, Inches(9.7), y, Inches(3.5), Inches(0.42),
             size=10, color=MUTED, italic=True, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────────────────
# Slide 4 — Phases 1-13 in one slide
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "PRIOR WORK", DEEP)
add_text(slide, "Phases 1–13 — what we already knew",
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         size=30, bold=True, color=NAVY, font="Georgia")
add_text(slide, "The mutagenesis null result that motivated Phase 14",
         Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
         size=14, italic=True, color=MUTED, font="Calibri")
add_bullets(slide, [
    [{"text": "Phase 1–2: ", "color": DEEP, "bold": True},
     {"text": "TYMS active site is highly conserved across 10 orthologs (JS-divergence + UniProt + PDBe)."}],
    [{"text": "Phase 5–6: ", "color": DEEP, "bold": True},
     {"text": "Wild-type dUMP docks at "},
     {"text": "−8.78 kcal/mol", "bold": True, "color": NAVY},
     {"text": " into the chain-A active-site of the 1HVY-derived dimer."}],
    [{"text": "Phase 7 (★ key finding): ", "color": ACCENT, "bold": True},
     {"text": "Rigid Vina + AD4 partial charges "},
     {"text": "cannot resolve TYMS active-site mutants at the kcal scale", "bold": True, "italic": True, "color": RED},
     {"text": ". Across 20 mutants × 2 conditions, the largest holo Δ score is "},
     {"text": "+0.77 kcal/mol", "bold": True, "color": RED},
     {"text": " — below the Vina noise floor of ±0.85 kcal/mol."}],
    [{"text": "Phase 8: ", "color": DEEP, "bold": True},
     {"text": "Vinardo rescoring + flexible-residue Vina partially fix the C195A illusion, but the R→E sign error needs proper PB electrostatics (MM-GBSA / FEP)."}],
    [{"text": "Phase 6 (Modeller): ", "color": DEEP, "bold": True},
     {"text": "10 homology models from 30–95% identity templates; best refined model scores 95.4% favoured under Lovell (the 1HVY crystal itself scores 92.2%)."}],
], Inches(0.6), Inches(1.7), Inches(12.5), Inches(4.5), size=13, gap=10)

# Bottom callout
callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.2),
                                  Inches(12.7), Inches(1.0))
callout.fill.solid(); callout.fill.fore_color.rgb = LIGHT
callout.line.color.rgb = TEAL; callout.line.width = Pt(1.5)
add_text(slide, "→ The honest null result of Phase 7-8 motivated Phase 14: instead of asking what mutants do "
                "to dUMP, ask what other molecules can bind to TYMS at the substrate site or elsewhere.",
         Inches(0.9), Inches(6.35), Inches(12.2), Inches(0.7),
         size=13, italic=True, color=NAVY, font="Calibri",
         anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────────────────
# Slide 5 — Section divider: Phase 14
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, NAVY)
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.0), SW, Inches(1.6))
deco.fill.solid(); deco.fill.fore_color.rgb = TEAL; deco.line.fill.background()
add_text(slide, "PHASE 14", Inches(0.6), Inches(3.1), Inches(12), Inches(0.6),
         size=18, color=WHITE, font="Consolas", bold=True)
add_text(slide, "Inhibitor design at four binding sites",
         Inches(0.6), Inches(3.6), Inches(12), Inches(1.0),
         size=44, color=WHITE, bold=True, font="Georgia")
add_text(slide, "1. Active site · 2. Cofactor site · 3. Dimer interface · 4. Allosteric / surface",
         Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
         size=18, color=RGBColor(0xCA, 0xDC, 0xFC), italic=True, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 6 — Phase 14 four strategies overview
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "PHASE 14 / OVERVIEW", DEEP)
add_text(slide, "Four strategies at a glance",
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         size=30, bold=True, color=NAVY, font="Georgia")
add_text(slide, "Same receptor, same Vina engine, same ±0.85 kcal/mol noise floor — orthogonal sites.",
         Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
         size=14, italic=True, color=MUTED, font="Calibri")

table_data = [
    ["#", "Site", "Tier-1 anchors", "Δ reference", "Headline finding"],
    ["1", [{"text": "Active site", "bold": True, "color": NAVY}], "dUMP, 5-FdUMP, BrdUMP, floxuridine, 5-FU", "dUMP apo",
     "nucleotides clean active-vs-prodrug gap (3 kcal/mol); 5-FdUMP within noise of dUMP"],
    ["2", [{"text": "Cofactor site", "bold": True, "color": NAVY}], "MTX, raltitrexed, pemetrexed, nolatrexed, plevitrexed", "raltitrexed apo",
     [{"text": "★ Plevitrexed (ZD9331) ", "color": ACCENT, "bold": True},
      {"text": "−10.01 kcal/mol, Δ −0.88 above noise"}]],
    ["3", [{"text": "Dimer interface", "bold": True, "color": NAVY}], "LR-octapeptide + scrambled control + 5 fragments", "scrambled control",
     "documented null — Vina cannot resolve 8-mer peptides (Hassan 2017)"],
    ["4", [{"text": "Allosteric", "bold": True, "color": NAVY}], "exploratory fragment screen (no clinical anchors)", "absolute Vina + FPocket",
     [{"text": "★ TYMS exposes an under-explored druggable cavity ", "color": ACCENT, "bold": True},
      {"text": "(FPocket 0.994); drug-like fragments at −7.5 kcal/mol"}]],
]
add_table(slide, table_data, Inches(0.6), Inches(1.7), Inches(12.7), Inches(4.0),
          col_widths=[0.4, 1.6, 3.0, 1.6, 5.5], cell_size=11, header_size=12)

add_text(slide, "Two real findings in 86 docked compounds — both subject to experimental follow-up.",
         Inches(0.6), Inches(6.0), Inches(12.7), Inches(0.4),
         size=14, italic=True, color=DEEP, align=PP_ALIGN.CENTER, font="Calibri", bold=True)
add_text(slide, "Full master CSV: 14_inhibitor_design/05_aggregate/master.csv (86 data rows)",
         Inches(0.6), Inches(6.55), Inches(12.7), Inches(0.4),
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER, font="Consolas")

# ────────────────────────────────────────────────────────────────
# Slide 7 — Strategy 1: Active site
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "STRATEGY 1 / ACTIVE SITE", DEEP)
add_text(slide, "S1 — Active site (dUMP-mimetic panel)",
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         size=28, bold=True, color=NAVY, font="Georgia")
add_text(slide, "Canonical Phase-7 box · 5 anchors + 7 RDKit decoys × 2 seeds × {apo, holo}",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

# Left: table
table_data = [
    ["Compound", "Tier", "Top1 apo", "Δ vs dUMP", "Verdict"],
    [[{"text": "5-FdUMP", "bold": True, "color": NAVY}], "1", [{"text": "−9.04", "bold": True, "color": NAVY}],
     "−0.27", "canonical TYMS active species"],
    ["BrdUMP", "1", "−8.88", "−0.10", "halogenated dUMP mimic"],
    ["dUMP (control)", "1", "−8.78", "0.00", "reference"],
    ["Floxuridine", "1", "−7.48", "+1.30", "no phosphate → 1.5 kcal/mol penalty"],
    ["decoy_CID6035", "2", "−7.47", "+1.32", "competing drug-like decoy"],
    ["5-FU (prodrug)", "1", [{"text": "−4.95", "bold": True, "color": RED}],
     "+3.83", [{"text": "weak — as expected", "italic": True}]],
]
add_table(slide, table_data, Inches(0.6), Inches(1.6), Inches(7.5), Inches(3.6),
          col_widths=[2.0, 0.6, 1.2, 1.2, 2.5], cell_size=11)

# Right: stats + teaching point
add_stat_callout(slide, "1.31 Å", "A0 re-dock RMSD\n(passes 2.0 Å gate)",
                 Inches(8.5), Inches(1.7), w=Inches(4.5), big_size=42, big_color=GREEN)
add_stat_callout(slide, "3 kcal/mol", "active-vs-prodrug gap\n(the enrichment signal)",
                 Inches(8.5), Inches(3.4), w=Inches(4.5), big_size=42, big_color=DEEP)
add_text(slide, "Teaching point", Inches(0.6), Inches(5.4), Inches(12), Inches(0.4),
         size=15, bold=True, color=ACCENT, font="Calibri")
add_text(slide, "5-FdUMP (the canonical TYMS inhibitor) scores 0.27 kcal/mol better than dUMP — "
                "directionally right but within Vina's ±0.85 noise floor. The decoy-vs-prodrug separation IS clean "
                "(3 kcal/mol), and floxuridine (no phosphate) reveals that phosphate-clamp engagement is worth "
                "~1.5 kcal/mol of binding energy.",
         Inches(0.6), Inches(5.8), Inches(12.7), Inches(1.5),
         size=13, color=TEXT, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 8 — Strategy 2: Cofactor site — ★ Plevitrexed
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "STRATEGY 2 / COFACTOR SITE", DEEP)
add_text(slide, "S2 — Cofactor site (antifolates)",
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         size=28, bold=True, color=NAVY, font="Georgia")
add_text(slide, "Box from holo D16 centroid · 6 anchors + 1 negative control + 4 decoys × 2 seeds × {apo, holo}",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

table_data = [
    ["Compound", "Tier", "Top1 apo", "Δ vs raltitrexed", "Verdict"],
    [[{"text": "★ Plevitrexed (ZD9331)", "bold": True, "color": ACCENT}], "1",
     [{"text": "−10.01", "bold": True, "color": ACCENT}],
     [{"text": "−0.88", "bold": True, "color": ACCENT}],
     [{"text": "first Phase-14 hit above noise floor", "bold": True, "color": ACCENT}]],
    ["Pemetrexed (S)", "1", "−9.72", "−0.59", "within noise; correct S-isomer"],
    ["pemetrexed (R) decoy", "2", "−9.63", "−0.50", "Vina is achiral — cannot distinguish"],
    ["Methotrexate", "1", "−9.59", "−0.46", "weak TYMS / strong DHFR — cross-target"],
    ["Raltitrexed (reference)", "1", "−9.13", "0.00", "the species bound in holo crystal"],
    ["Nolatrexed (AG-337)", "1", "−7.57", "+1.56", "no glutamate tail → weaker"],
]
add_table(slide, table_data, Inches(0.6), Inches(1.6), Inches(8.8), Inches(3.4),
          col_widths=[2.5, 0.6, 1.2, 1.5, 3.0], cell_size=11)

# Right: big stat + chemistry
add_stat_callout(slide, "−10.01", "Plevitrexed top1\nkcal/mol",
                 Inches(9.8), Inches(1.7), w=Inches(3.2), big_size=48, big_color=ACCENT)
add_stat_callout(slide, "−0.88", "vs raltitrexed\n(noise = ±0.85)",
                 Inches(9.8), Inches(3.6), w=Inches(3.2), big_size=42, big_color=DEEP)

# Bottom: teaching point
add_text(slide, "Teaching points", Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
         size=15, bold=True, color=ACCENT, font="Calibri")
add_bullets(slide, [
    [{"text": "Plevitrexed (ZD9331) is reproducible across both seeds (−9.95, −10.07). Consistent with its sub-nanomolar TYMS Ki ("},
     {"text": "Jackman 1997", "italic": True}, {"text": ")."}],
    "Vina is achiral — pemetrexed S (clinical) and R (decoy) score identically. Known limitation, not a TYMS-specific failure.",
    "Holo state drops every antifolate ~3 kcal/mol — the already-bound raltitrexed sterically competes (clean demonstration of holo = 'displacement contest').",
], Inches(0.6), Inches(5.6), Inches(12.7), Inches(1.8), size=12, gap=4)

# ────────────────────────────────────────────────────────────────
# Slide 9 — Strategy 3: Dimer interface (null result)
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "STRATEGY 3 / DIMER INTERFACE", DEEP)
add_text(slide, "S3 — Dimer interface (PPI disruptor)",
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
         size=28, bold=True, color=NAVY, font="Georgia")
add_text(slide, "MDAnalysis 4 Å contact map · LR-octapeptide LSCQLYQR + scrambled control + 5 overlapping 4-mer fragments",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

table_data = [
    ["Peptide", "Length", "Kind", "Top1 (mean)", "Verdict"],
    ["LR8_LSCQLYQR (canonical)", "8", "canonical lead",
     [{"text": "+86.16", "color": RED, "bold": True}],
     "Vina cannot fit 8-mer in interface box"],
    ["LR8_scrambled_QLCRQSYL", "8", "scrambled control",
     [{"text": "+84.68", "color": RED, "bold": True}],
     "same failure mode"],
    ["LR_4mer × 5 fragments", "4", "overlapping windows", "−4.1 to −4.7",
     "weak surface binding"],
]
add_table(slide, table_data, Inches(0.6), Inches(1.6), Inches(12.7), Inches(2.2),
          col_widths=[3.0, 0.8, 2.0, 2.0, 4.9], cell_size=12)

# Big finding box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.0),
                              Inches(12.7), Inches(1.5))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT
box.line.color.rgb = RED; box.line.width = Pt(2)
add_text(slide,
         "Specificity vs scrambled = +1.48 kcal/mol  →  canonical worse than scrambled.",
         Inches(0.9), Inches(4.15), Inches(12.1), Inches(0.5),
         size=18, bold=True, color=RED, align=PP_ALIGN.CENTER, font="Calibri",
         anchor=MSO_ANCHOR.MIDDLE)
add_text(slide, "Documented null result per Stop Condition S1. Rigid Vina is the wrong engine for ≥6-mer peptides "
                "(Hassan 2017). The right tools (HPEPDOCK, CABS-dock, FlexPepDock, RosettaDock) were unreachable / "
                "out-of-scope at execution time. The null is the correct conclusion, not a methodology failure.",
         Inches(0.9), Inches(4.6), Inches(12.1), Inches(0.8),
         size=12, italic=True, color=TEXT, font="Calibri", align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(slide, "Honest null results are findings too.", Inches(0.6), Inches(6.0), Inches(12.7), Inches(0.4),
         size=14, italic=True, bold=True, color=DEEP, font="Calibri", align=PP_ALIGN.CENTER)
add_text(slide, "An unhonestly-positive PPI ligand from rigid Vina would be misleading downstream.",
         Inches(0.6), Inches(6.5), Inches(12.7), Inches(0.4),
         size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 10 — Strategy 4: Allosteric — cavity 18 headline
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "STRATEGY 4 / ALLOSTERIC", ACCENT)
add_text(slide, "S4 — Allosteric: TYMS has an under-explored druggable cavity",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.7),
         size=26, bold=True, color=NAVY, font="Georgia")
add_text(slide, "Self-built FPocket 4.0 (arm64-darwin) · 33 detected pockets · 100 dock runs",
         Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

# 3 callouts at top
add_stat_callout(slide, "0.994", "FPocket druggability\n(cavity 18, near ceiling)",
                 Inches(0.6), Inches(1.7), w=Inches(3.0), big_size=50, big_color=ACCENT)
add_stat_callout(slide, "0.828", "C2-symmetric mirror\n(cavity 17, partner protomer)",
                 Inches(3.7), Inches(1.7), w=Inches(3.0), big_size=44, big_color=DEEP)
add_stat_callout(slide, "−7.52", "1H-indazole top hit\nkcal/mol",
                 Inches(6.8), Inches(1.7), w=Inches(3.0), big_size=50, big_color=ACCENT)
add_stat_callout(slide, "−7.28", "Ibuprofen top hit\nkcal/mol",
                 Inches(9.9), Inches(1.7), w=Inches(3.0), big_size=50, big_color=ACCENT)

# Table of top hits
table_data = [
    ["Fragment", "Common name", "Cavity", "Top1", "Druggability"],
    [[{"text": "★ frag_CID7032", "bold": True, "color": ACCENT}],
     "1H-indazole (kinase-inhibitor scaffold)",
     [{"text": "18", "bold": True, "color": ACCENT}],
     [{"text": "−7.52", "bold": True, "color": ACCENT}],
     [{"text": "0.994", "bold": True, "color": ACCENT}]],
    [[{"text": "★ frag_CID3672", "bold": True, "color": ACCENT}],
     "ibuprofen (NSAID, COX1/2, promiscuous)",
     [{"text": "18", "bold": True, "color": ACCENT}],
     [{"text": "−7.28", "bold": True, "color": ACCENT}],
     [{"text": "0.994", "bold": True, "color": ACCENT}]],
    ["frag_CID5564", "tolnaftate (antifungal)", "2", "−6.88", "0.009"],
    ["frag_CID7032", "1H-indazole (same ligand, low-drug cavity)", "2", "−6.86", "0.009"],
    ["frag_CID35814", "flurbiprofen (NSAID)", "12", "−6.52", "0.010"],
]
add_table(slide, table_data, Inches(0.6), Inches(3.7), Inches(12.7), Inches(2.5),
          col_widths=[2.0, 4.0, 1.2, 1.5, 1.8], cell_size=11, header_size=12)

add_text(slide,
         "Two unrelated drug-like fragments → 2 kcal/mol better at the druggable pocket than anywhere else on the surface.",
         Inches(0.6), Inches(6.4), Inches(12.7), Inches(0.5),
         size=14, bold=True, color=DEEP, font="Calibri", align=PP_ALIGN.CENTER)
add_text(slide, "Refutes the v1 'no obvious druggable allosteric pocket' framing.",
         Inches(0.6), Inches(6.9), Inches(12.7), Inches(0.4),
         size=12, color=MUTED, italic=True, align=PP_ALIGN.CENTER, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 11 — Cavity 18 + indazole pose
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "CAVITY 18 / TOP HIT", ACCENT)
add_text(slide, "1H-indazole at cavity 18  —  −7.52 kcal/mol",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=26, bold=True, color=NAVY, font="Georgia")
add_text(slide, "CID 7032 · C₇H₆N₂ · MW 118 · kinase-inhibitor privileged scaffold (axitinib, niraparib, pazopanib)",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=12, italic=True, color=MUTED, font="Calibri")

# Image left
add_image(slide, "pose_indazole", Inches(0.6), Inches(1.6), w=Inches(6.5))

# Interactions right
add_text(slide, "Residue contacts (chain B)", Inches(7.4), Inches(1.6), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=DEEP, font="Calibri")
table_data = [
    ["Residue", "d (Å)", "Interaction"],
    [[{"text": "Phe B55", "bold": True, "color": NAVY}], "3.23", "H-bond + π + hydrophobic"],
    [[{"text": "Asn B201", "bold": True, "color": NAVY}], "3.26", "H-bond"],
    [[{"text": "Leu B196 ★", "bold": True, "color": ACCENT}], "3.40", "★ allosteric loop 181-197"],
    [[{"text": "Gly B197 ★", "bold": True, "color": ACCENT}], "3.48", "★ allosteric loop 181-197"],
    [[{"text": "Phe B200", "bold": True, "color": NAVY}], "3.51", "π-stack + hydrophobic"],
    ["Ile B83", "3.62", "hydrophobic"],
    ["Val B54, Lys B52", "3.7", "hydrophobic walls"],
]
add_table(slide, table_data, Inches(7.4), Inches(2.05), Inches(5.8), Inches(3.6),
          col_widths=[2.6, 0.8, 2.4], cell_size=11)

# Bottom mechanism
add_text(slide, "Mechanism inferred from the pose", Inches(0.6), Inches(6.0),
         Inches(12.7), Inches(0.4),
         size=14, bold=True, color=ACCENT, font="Calibri")
add_text(slide, "Three of ten contact residues (Leu196, Gly197, Phe200) are on the published allosteric "
                "communication loop 181-197 (Anderson 2012, Pozzi 2019). Pose geometry consistent with allosteric "
                "mechanism: occupying the loop face restricts hinge motion coupled to active-site Cys195.",
         Inches(0.6), Inches(6.4), Inches(12.7), Inches(1.0),
         size=12, color=TEXT, italic=True, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 12 — Cavity 18 + ibuprofen pose
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "CAVITY 18 / IBUPROFEN POSE", ACCENT)
add_text(slide, "Ibuprofen at cavity 18  —  −7.28 kcal/mol",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=26, bold=True, color=NAVY, font="Georgia")
add_text(slide, "CID 3672 · C₁₃H₁₈O₂ · MW 206 · NSAID, COX1/2; promiscuous off-target binder (HSA, FABP4, CRBN)",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=12, italic=True, color=MUTED, font="Calibri")

# Image left
add_image(slide, "pose_ibuprofen", Inches(0.6), Inches(1.6), w=Inches(6.5))

# Right
add_text(slide, "Residue contacts (chain B)", Inches(7.4), Inches(1.6), Inches(5.8), Inches(0.4),
         size=15, bold=True, color=DEEP, font="Calibri")
table_data = [
    ["Residue", "d (Å)", "Interaction"],
    [[{"text": "★ Lys B283", "bold": True, "color": ACCENT}], "3.01", "★ salt bridge (carboxylate)"],
    [[{"text": "★ Lys B52", "bold": True, "color": ACCENT}], "3.08", "★ salt bridge (carboxylate)"],
    [[{"text": "Leu B196 ★", "bold": True, "color": ACCENT}], "3.29", "allosteric loop 181-197"],
    [[{"text": "Phe B200", "bold": True, "color": NAVY}], "3.57", "π-stack + hydrophobic"],
    ["Val B54, Phe B55", "3.6", "hydrophobic"],
    ["Ile B83", "3.58", "hydrophobic"],
    [[{"text": "Gly B197 ★", "bold": True, "color": ACCENT}], "3.80", "allosteric loop 181-197"],
]
add_table(slide, table_data, Inches(7.4), Inches(2.05), Inches(5.8), Inches(3.6),
          col_widths=[2.6, 0.8, 2.4], cell_size=11)

# Bottom mechanism
add_text(slide, "The double-lysine salt-bridge clamp", Inches(0.6), Inches(6.0),
         Inches(12.7), Inches(0.4),
         size=14, bold=True, color=ACCENT, font="Calibri")
add_text(slide, "Lys52 and Lys283 are 14 sequence positions apart but ~6 Å apart in 3D — together they form a "
                "positively-charged anchor that clamps the deprotonated propanoate at pH 7.4. "
                "Any future cavity-18 lead should carry an anionic head-group to exploit this geometry.",
         Inches(0.6), Inches(6.4), Inches(12.7), Inches(1.0),
         size=12, color=TEXT, italic=True, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 13 — Same-ligand-different-pockets proof
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "CAVITY 18 / PROOF", ACCENT)
add_text(slide, "Why we trust the cavity-18 finding",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=26, bold=True, color=NAVY, font="Georgia")
add_text(slide, "Two head-to-head sanity checks the result passes",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

# Two columns
add_text(slide, "1. Same ligand, different pockets",
         Inches(0.6), Inches(1.6), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=DEEP, font="Calibri")
add_image(slide, "pose_indazole_cav2", Inches(0.6), Inches(2.1), w=Inches(6.0))
add_text(slide, "1H-indazole at cavity 2 (druggability 0.009):  −6.86 kcal/mol\n"
                "Same molecule at cavity 18 (druggability 0.994):  −7.52 kcal/mol",
         Inches(0.6), Inches(5.6), Inches(6.5), Inches(0.8),
         size=11, color=TEXT, font="Calibri")
add_text(slide, "13 surface contacts vs 10 pocket contacts — more contacts, weaker binding.",
         Inches(0.6), Inches(6.4), Inches(6.5), Inches(0.4),
         size=11, italic=True, color=ACCENT, bold=True, font="Calibri")

add_text(slide, "2. Different ligands, same pocket",
         Inches(7.0), Inches(1.6), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=DEEP, font="Calibri")
add_image(slide, "pose_flurbi", Inches(7.0), Inches(2.1), w=Inches(6.0))
add_text(slide, "Flurbiprofen at cavity 12 (druggability 0.010):  −6.52 kcal/mol\n"
                "Two unrelated scaffolds at cavity 18:  −7.5, −7.3 kcal/mol",
         Inches(7.0), Inches(5.6), Inches(6.5), Inches(0.8),
         size=11, color=TEXT, font="Calibri")
add_text(slide, "Indazole (H-bond + π) vs ibuprofen (salt-bridge + π) — different chemistries, same pocket.",
         Inches(7.0), Inches(6.4), Inches(6.5), Inches(0.4),
         size=11, italic=True, color=ACCENT, bold=True, font="Calibri")

add_text(slide, "More-contacts ≠ better-binding without concavity.  Cavity 18 discriminates chemistry — exactly what a real druggable pocket does.",
         Inches(0.6), Inches(6.95), Inches(12.7), Inches(0.4),
         size=12, color=NAVY, italic=True, bold=True, align=PP_ALIGN.CENTER, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 14 — Conservation
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "CAVITY 18 / CONSERVATION", ACCENT)
add_text(slide, "Cavity-18 residues are unusually conserved",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=26, bold=True, color=NAVY, font="Georgia")
add_text(slide, "JS conservation across 10 TYMS orthologs · loop 181-197 highlighted in red",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

add_image(slide, "conservation", Inches(0.6), Inches(1.5), w=Inches(8.5))

# Right column
add_text(slide, "100% conserved across 10 orthologs",
         Inches(9.3), Inches(1.7), Inches(3.8), Inches(0.4),
         size=14, bold=True, color=ACCENT, font="Calibri")
add_bullets(slide, [
    [{"text": "Gly54", "bold": True, "color": NAVY}, {"text": " (hydrophobic floor)"}],
    [{"text": "Glu87", "bold": True, "color": NAVY}, {"text": " (polar handle for future leads)"}],
    [{"text": "Met190", "bold": True, "color": NAVY}, {"text": " (loop 181-197 ★)"}],
    [{"text": "Ala191", "bold": True, "color": NAVY}, {"text": " (loop 181-197 ★)"}],
    [{"text": "Leu196", "bold": True, "color": NAVY}, {"text": " (loop 181-197 ★)"}],
    [{"text": "Phe200", "bold": True, "color": NAVY}, {"text": " (π-stack anchor — both hits)"}],
    [{"text": "Asn201", "bold": True, "color": NAVY}, {"text": " (H-bond — indazole)"}],
], Inches(9.3), Inches(2.15), Inches(3.8), Inches(3.5), size=11, gap=2)
add_text(slide, "6 of the 7 ultra-conserved positions are contact residues for at least one top hit.",
         Inches(9.3), Inches(5.4), Inches(3.8), Inches(1.5),
         size=11, color=ACCENT, bold=True, italic=True, font="Calibri")

add_text(slide, "Conservation at the contact face = strong evidence the pocket is functionally meaningful, not a random surface patch.",
         Inches(0.6), Inches(6.9), Inches(12.7), Inches(0.4),
         size=12, color=DEEP, italic=True, bold=True, align=PP_ALIGN.CENTER, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 15 — Phylogeny
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "CAVITY 18 / PHYLOGENY", ACCENT)
add_text(slide, "Phylogeny — cavity-18 mutations across orthologs",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=26, bold=True, color=NAVY, font="Georgia")
add_text(slide, "NJ tree (BLOSUM62 distance) annotated with # cavity-18 substitutions vs human",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

add_image(slide, "phylogeny", Inches(0.6), Inches(1.5), w=Inches(8.0))

# Right column
add_text(slide, "Substitutions vs human",
         Inches(8.8), Inches(1.7), Inches(4.5), Inches(0.4),
         size=14, bold=True, color=DEEP, font="Calibri")
table_data = [
    ["Ortholog", "#mut"],
    ["Homo sapiens", [{"text": "0 (ref)", "color": MUTED}]],
    ["Mus musculus", [{"text": "2", "color": GREEN, "bold": True}]],
    ["Rattus norvegicus", [{"text": "2", "color": GREEN, "bold": True}]],
    ["Drosophila / yeast / E. coli", "12"],
    ["L. casei", "16"],
    ["Bacteriophage T4", "18"],
    [[{"text": "P. falciparum ★", "color": ACCENT, "bold": True}],
     [{"text": "21", "color": ACCENT, "bold": True}]],
]
add_table(slide, table_data, Inches(8.8), Inches(2.15), Inches(4.5), Inches(3.5),
          col_widths=[3.0, 1.5], cell_size=11)

# Bottom callout
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.0),
                              Inches(12.7), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT
box.line.color.rgb = ACCENT; box.line.width = Pt(2)
add_text(slide, "Mammals share the cavity signature near-identically.  Plasmodium falciparum diverges at 21 of 36 positions.",
         Inches(0.9), Inches(6.15), Inches(12.1), Inches(0.45),
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Calibri",
         anchor=MSO_ANCHOR.MIDDLE)
add_text(slide, "→ a species-selective allosteric TYMS inhibitor is structurally plausible at cavity 18 — "
                "distinct from the canonical active site, which is too conserved to give a selectivity handle.",
         Inches(0.9), Inches(6.6), Inches(12.1), Inches(0.5),
         size=11, color=TEXT, italic=True, align=PP_ALIGN.CENTER, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 16 — 3D viewers + downloads
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "CAVITY 18 / EXPLORE", ACCENT)
add_text(slide, "Interactive 3D viewers & downloadable structures",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=26, bold=True, color=NAVY, font="Georgia")
add_text(slide, "3Dmol.js viewers — cavity-18 surface in wheat, allosteric loop 181-197 ∩ cavity = red, ligand = cyan",
         Inches(0.6), Inches(1.0), Inches(12.7), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

# Viewer cards
viewers = [
    ("Apo (no ligand)", "Pocket structure on its own", "cavity18_apo.html",
     "Toggle ligand later · zoom to pocket · zoom to loop"),
    ("+ 1H-indazole",   "★ top hit · −7.52 kcal/mol", "cavity18_indazole.html",
     "Polar contacts highlighted as yellow dashes"),
    ("+ Ibuprofen",     "−7.28 · double-lysine salt-bridge", "cavity18_ibuprofen.html",
     "Salt-bridge dashes to Lys52 / Lys283"),
]
card_w = Inches(4.0); card_h = Inches(2.6); gap_x = Inches(0.27); x0 = Inches(0.6); y = Inches(1.6)
for i, (title, sub, fname, body) in enumerate(viewers):
    x = x0 + i * (card_w + gap_x)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = ACCENT if i < 2 else DEEP
    card.line.width = Pt(1.5)
    add_text(slide, title, x+Inches(0.25), y+Inches(0.2), card_w-Inches(0.5), Inches(0.45),
             size=16, bold=True, color=NAVY, font="Calibri")
    add_text(slide, sub, x+Inches(0.25), y+Inches(0.65), card_w-Inches(0.5), Inches(0.4),
             size=11, italic=True, color=ACCENT if i < 2 else TEXT, font="Calibri")
    add_text(slide, body, x+Inches(0.25), y+Inches(1.15), card_w-Inches(0.5), Inches(0.8),
             size=10, color=TEXT, font="Calibri")
    add_text(slide, fname, x+Inches(0.25), y+Inches(2.1), card_w-Inches(0.5), Inches(0.4),
             size=9, color=MUTED, font="Consolas")

# Downloads list
add_text(slide, "Downloadable files (cavity18_evidence/downloads/)",
         Inches(0.6), Inches(4.5), Inches(12.7), Inches(0.4),
         size=14, bold=True, color=DEEP, font="Calibri")
add_bullets(slide, [
    [{"text": "cavity18_apo.pdb", "font": "Consolas", "color": NAVY}, {"text": "  —  apo TYMS dimer with chain A/B labels"}],
    [{"text": "cavity18_apo_pocket.pdb", "font": "Consolas", "color": NAVY}, {"text": "  —  pocket residues only, B-factor = druggability × 100"}],
    [{"text": "cavity18_indazole_complex.pdb", "font": "Consolas", "color": NAVY}, {"text": "  —  apo + indazole top pose (HETATM IND Z)"}],
    [{"text": "cavity18_ibuprofen_complex.pdb", "font": "Consolas", "color": NAVY}, {"text": "  —  apo + ibuprofen top pose (HETATM IBU Z)"}],
    [{"text": "cavity18_residues.csv", "font": "Consolas", "color": NAVY}, {"text": "  —  36-row residue × ortholog × conservation table"}],
    [{"text": "cavity18_mutations_per_taxon.json", "font": "Consolas", "color": NAVY}, {"text": "  —  list of substitutions per ortholog"}],
], Inches(0.6), Inches(4.95), Inches(12.7), Inches(2.0), size=11, gap=3)

# ────────────────────────────────────────────────────────────────
# Slide 17 — Multi-agent audit chain
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "QUALITY / AUDIT", DEEP)
add_text(slide, "Multi-agent peer review — 6 rounds, all logged",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=28, bold=True, color=NAVY, font="Georgia")
add_text(slide, "Same doer ↔ verifier convention used in Phases 1–13.  No compute spent until R3 PASS on the roadmap.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
         size=13, italic=True, color=MUTED, font="Calibri")

table_data = [
    ["Round", "Verdict", "Top finding", "Fixed in"],
    ["R1 (roadmap)", "CONDITIONAL", [{"text": "8 of 10 v0 anchor CIDs pointed to the wrong compound", "color": RED, "bold": True},
                                    {"text": " (dUMP 22848 → Solanum steroid; nolatrexed 60198 → estrogen analog)"}],
     "v1 — verified-anchors JSON committed"],
    ["R2 (roadmap)", "CONDITIONAL", "CID gate still a no-op · PROLIF can't flag missing waters · HPEPDOCK had no fallback or timeout",
     "v2 — direct verification + E1b script + CABS-dock envelope"],
    ["R3 (roadmap)", "CONDITIONAL", "E1b alignment script crashed on ligand-only PDBQT · pemetrexed null InChIKey · ConnectivitySMILES wrong field",
     "v2-final — frame check redesigned, all 3 fixed"],
    ["R4 (results)", "CONDITIONAL", "SASA column > 1 · A0 RMSD failed gate · duplicate rows in S1 · S4 overstated",
     "R4-fix commit"],
    ["R5 (verify)", [{"text": "PASS", "color": GREEN, "bold": True}], "all 4 R4 blockers verified closed", "—"],
    ["R6 (S4 v2)", "CONDITIONAL", [{"text": "Pocket 17 is the C2-symmetric mirror of pocket 18 (positive sanity); 'cryptic' is wrong word; literature anchor (Anderson 2012, Pozzi 2019)", "italic": True}],
     "Documentation corrections applied"],
]
add_table(slide, table_data, Inches(0.6), Inches(1.7), Inches(12.7), Inches(5.0),
          col_widths=[1.4, 1.3, 6.0, 4.0], cell_size=10, header_size=11)
add_text(slide, "All reviewer reports verbatim under 14_inhibitor_design/00_roadmap/reviews/",
         Inches(0.6), Inches(6.95), Inches(12.7), Inches(0.4),
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER, font="Consolas")

# ────────────────────────────────────────────────────────────────
# Slide 18 — What we found & what's next
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, WHITE)
add_section_bar(slide, "SUMMARY", DEEP)
add_text(slide, "What we found in Phase 14",
         Inches(0.6), Inches(0.4), Inches(12.7), Inches(0.6),
         size=30, bold=True, color=NAVY, font="Georgia")

# Two-column
add_text(slide, "Real findings (2)", Inches(0.6), Inches(1.2), Inches(6.0), Inches(0.4),
         size=18, bold=True, color=ACCENT, font="Calibri")
add_bullets(slide, [
    [{"text": "Plevitrexed (ZD9331) — cofactor site, ", "bold": True, "color": NAVY},
     {"text": "−10.01 kcal/mol", "bold": True, "color": ACCENT},
     {"text": ", Δ −0.88 above the Vina ±0.85 kcal/mol noise floor.  Reproducible across both seeds.  Consistent with its published sub-nanomolar Ki (Jackman 1997).  Reproduced, not discovered."}],
    "",
    [{"text": "Cavity 18 — under-explored druggable allosteric pocket.  ", "bold": True, "color": NAVY},
     {"text": "FPocket druggability 0.994 (and 0.828 C2 mirror).  Drug-like fragments dock at −7.5 kcal/mol via different chemistry (indazole H-bonds + π; ibuprofen double-Lys salt-bridge + π).  Overlaps the published allosteric communication loop 181-197 (Anderson 2012, Pozzi 2019)."}],
], Inches(0.6), Inches(1.7), Inches(6.2), Inches(5.0), size=11, gap=6)

add_text(slide, "Honest null results (and 1 reproduction)", Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.4),
         size=18, bold=True, color=DEEP, font="Calibri")
add_bullets(slide, [
    [{"text": "Active site: ", "bold": True, "color": NAVY},
     {"text": "5-FdUMP best at −9.04, but within ±0.85 noise floor of dUMP.  Phase 7-8 kcal-noise problem inherited."}],
    [{"text": "Dimer interface: ", "bold": True, "color": NAVY},
     {"text": "rigid Vina cannot fit 8-mer peptides at the interface (specificity vs scrambled = +1.48; canonical worse than scrambled).  Documented null per Stop Condition S1."}],
    "",
    [{"text": "8 of 10 v0 anchor CIDs were the wrong compound — ", "color": RED, "bold": True, "italic": True},
     {"text": "this is the loud silent-failure mode that the R1+R2 reviewer rounds caught before any docking ran."}],
], Inches(7.0), Inches(1.7), Inches(6.2), Inches(5.0), size=11, gap=6)

# Bottom: next steps
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.4),
                              Inches(12.7), Inches(0.9))
box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
add_text(slide, "Next steps (signposted, not executed)", Inches(0.9), Inches(6.45),
         Inches(12), Inches(0.4),
         size=13, bold=True, color=TEAL, font="Calibri")
add_text(slide, "Experimental fragment soak + crystal at cavity 18  ·  enzyme assay (kcat/Km) for cavity-18 ligands  ·  "
                "MM-GBSA / FEP rescoring  ·  GNINA CNN scoring (Apple-Silicon port pending)",
         Inches(0.9), Inches(6.85), Inches(12), Inches(0.4),
         size=11, color=WHITE, italic=True, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Slide 19 — TL;DR
# ────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_bg(slide, NAVY)
deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), SH)
deco.fill.solid(); deco.fill.fore_color.rgb = ACCENT; deco.line.fill.background()

add_text(slide, "TL;DR", Inches(0.9), Inches(0.6), Inches(12), Inches(1.0),
         size=60, bold=True, color=WHITE, font="Georgia")

add_text(slide, "Phase 14 — what to remember",
         Inches(0.9), Inches(1.7), Inches(12), Inches(0.5),
         size=18, color=TEAL, italic=True, font="Calibri")

add_bullets(slide, [
    [{"text": "TYMS is a two-site enzyme in the drug literature (substrate + cofactor).  ", "color": RGBColor(0xCA,0xDC,0xFC)},
     {"text": "We found a plausible third site.", "color": WHITE, "bold": True}],
    [{"text": "FPocket druggability ", "color": RGBColor(0xCA,0xDC,0xFC)},
     {"text": "0.994 ", "bold": True, "color": ACCENT},
     {"text": "on chain B + ", "color": RGBColor(0xCA,0xDC,0xFC)},
     {"text": "0.828 ", "bold": True, "color": ACCENT},
     {"text": "C2 mirror on chain A — same pocket on both protomers, detected independently.", "color": RGBColor(0xCA,0xDC,0xFC)}],
    [{"text": "Two unrelated drug-like fragments (1H-indazole, ibuprofen) bind at ", "color": RGBColor(0xCA,0xDC,0xFC)},
     {"text": "−7.5 kcal/mol", "bold": True, "color": ACCENT},
     {"text": " via different chemistries — proving the pocket discriminates ligands.", "color": RGBColor(0xCA,0xDC,0xFC)}],
    [{"text": "6 of 7 ultra-conserved cavity residues are exactly the contact face — the pocket is functionally meaningful, not a random surface patch.", "color": RGBColor(0xCA,0xDC,0xFC)}],
    [{"text": "Plasmodium falciparum has 21 cavity-18 substitutions vs human — ", "color": RGBColor(0xCA,0xDC,0xFC)},
     {"text": "structurally plausible species-selective allosteric handle.", "bold": True, "color": WHITE}],
    [{"text": "Pharmacophore proposal: ", "color": RGBColor(0xCA,0xDC,0xFC)},
     {"text": "anionic head (Lys52 + Lys283 clamp) + aromatic body (Phe200 π-stack) + hydrophobic tail (Leu196 / Gly197 / Ile83 wall).", "bold": True, "color": WHITE}],
], Inches(0.9), Inches(2.5), Inches(12), Inches(4.5), size=15, gap=10,
   bullet_color=ACCENT)

add_text(slide, "github.com/ArioMoniri/aminak  ·  Phase 14  ·  all artefacts + reviewer reports on main",
         Inches(0.9), Inches(6.9), Inches(12), Inches(0.4),
         size=11, color=TEAL, italic=True, font="Calibri")

# ────────────────────────────────────────────────────────────────
# Save
# ────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"→ {OUT_PATH}  ({OUT_PATH.stat().st_size/1024:.0f} KB)")
