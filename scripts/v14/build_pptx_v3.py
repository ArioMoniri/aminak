#!/usr/bin/env python3
"""Phase 14 — v3 PPTX with new slides + deep-link hyperlinks + ZIP bundle.

Fixes/additions vs v2 per user feedback:
  - Corrected WT statement: WT_apo = −8.55 (Phase 7 multi-replica mean), not −8.78
    (which was the deterministic single-seed value); show both honestly
  - Cleaner dimer-aware images (Phase A1 outputs)
  - New conservation-page layout (Phase A2 — JS plot + mutation table + phylogeny)
  - Smina full Phase-7 panel comparison (Phase A3)
  - Modeller vs AlphaFold backbone+geometry comparison (Phase A4)
  - PPI / BSA dimer-interface analysis (Phase A5)
  - Ramachandran improvement before/after (Phase A6)
  - MM-GBSA roadmap slide (Phase A7)
  - HADDOCK3 roadmap slide (Phase A8)
  - "Plasmodium selectivity explained" slide (the question the user asked)
  - Deep-link hyperlinks on key text → file:// URLs to local PDB/HTML/CSV

Output: 14_inhibitor_design/presentation/aminak_phase14_summary.pptx
        Also copied to /Users/ario/Downloads/Aminak Phase 14 Summary.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import shutil

# Palette — same as v2 (charcoal + warm beige + coral accent)
INK     = RGBColor(0x14, 0x18, 0x1F)
DARK    = RGBColor(0x21, 0x29, 0x33)
CHAR    = RGBColor(0x3A, 0x42, 0x4E)
MUTE    = RGBColor(0x8A, 0x94, 0xA1)
BONE    = RGBColor(0xF7, 0xF4, 0xEE)
PAPER   = RGBColor(0xFD, 0xFB, 0xF7)
RULE    = RGBColor(0xD8, 0xD2, 0xC6)
ACCENT  = RGBColor(0xE6, 0x55, 0x3F)
GREEN   = RGBColor(0x3E, 0x7B, 0x4F)
LINK    = RGBColor(0x06, 0x5A, 0x82)   # deep blue — for hyperlinks
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Inter SemiBold"
BODY = "Inter"
MONO = "JetBrains Mono"

REPO = Path(__file__).resolve().parents[2]
OUT_REPO = REPO / "14_inhibitor_design" / "presentation" / "aminak_phase14_summary.pptx"
OUT_USER = Path("/Users/ario/Downloads/Aminak Phase 14 Summary.pptx")
FIG = REPO / "14_inhibitor_design" / "presentation" / "figures"

# Image inventory — including the new clean renders + analysis batch outputs
IMG = {
    "workflow":     REPO / "workflow_diagram_v3.png",
    "conservation_multipanel": FIG / "conservation_multipanel.png",
    "dimer_overview_clean": FIG / "dimer_overview_clean.png",
    "dimer_activesite_clean": FIG / "dimer_activesite_clean.png",
    "cavity18_carve_clean": FIG / "cavity18_carve_clean.png",
    "holo_dump_cofactor": FIG / "holo_dump_cofactor_clean.png",
    "delta_apo":    REPO / "08d_analysis_v4" / "delta_vina_apo_holo.png",
    "phase8_flex":  REPO / "13_phase8" / "02_flexres" / "flex_vs_rigid.png",
    "phase8_alt":   REPO / "13_phase8" / "01_alt_scoring" / "alt_scoring_compare.png",
    "rama_best":    REPO / "10b_modeller_refined" / "04_refined_lovell" / "ramachandran_lovell_best_refined.png",
    "rama_compare": FIG / "ramachandran_before_after.png",
    "modeller_af":  FIG / "modeller_vs_alphafold.png",
    "ppi_dimer":    FIG / "ppi_dimer_interface.png",
    "smina_full":   FIG / "smina_full_panel.png",
    "openmm_gb":    REPO / "14_inhibitor_design" / "07_advanced_methods" / "openmm_gb_rescore" / "openmm_gb_plot.png",
    "smina_rescore": REPO / "14_inhibitor_design" / "06_smina_rescore" / "rescore_plot.png",
    "phase14_dist": REPO / "14_inhibitor_design" / "figures" / "fig1_distributions.png",
    "phase14_delta":REPO / "14_inhibitor_design" / "figures" / "fig2_delta_ranking.png",
    "phase14_apoholo": REPO / "14_inhibitor_design" / "figures" / "fig3_apo_holo_gap.png",
    "phase14_tier": REPO / "14_inhibitor_design" / "figures" / "fig4_tier_separation.png",
    "pose_indazole":  REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav18_CID7032.png",
    "pose_ibuprofen": REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav18_CID3672.png",
    "pose_indazole_cav2": REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav2_CID7032.png",
    "pose_flurbi":    REPO / "14_inhibitor_design" / "04_allosteric" / "poses" / "cav12_CID35814.png",
    "conservation_cav18": REPO / "14_inhibitor_design" / "04_allosteric" / "cavity18_evidence" / "figures" / "cavity18_conservation.png",
    "phylogeny_cav18":    REPO / "14_inhibitor_design" / "04_allosteric" / "cavity18_evidence" / "figures" / "cavity18_phylogeny_annot.png",
}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def text(slide, txt, x, y, w, h, size=14, bold=False, color=CHAR, align=PP_ALIGN.LEFT,
         font=BODY, italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    # normalise: txt can be (a) a string  (b) a list of strings  (c) a list of dicts (single para, multi-run)
    #                       (d) a list of lists (each inner list = one paragraph with multi-run dicts)
    if isinstance(txt, list):
        if txt and isinstance(txt[0], dict):
            paragraphs = [txt]  # single paragraph, multi-run
        else:
            paragraphs = txt    # multiple paragraphs
    else:
        paragraphs = [txt]
    for i, line in enumerate(paragraphs):
        if i > 0:
            p = tf.add_paragraph(); p.alignment = align
        if isinstance(line, list):
            for run_spec in line:
                r = p.add_run()
                r.text = run_spec["text"]
                r.font.size = Pt(run_spec.get("size", size))
                r.font.bold = run_spec.get("bold", bold)
                r.font.italic = run_spec.get("italic", italic)
                r.font.color.rgb = run_spec.get("color", color)
                r.font.name = run_spec.get("font", font)
                # hyperlink support
                if "hyperlink" in run_spec:
                    r.hyperlink.address = run_spec["hyperlink"]
        else:
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.bold = bold
            r.font.italic = italic; r.font.color.rgb = color
            r.font.name = font
    return tb

def add_image(slide, key, x, y, w=None, h=None):
    p = IMG[key]
    if not p.exists():
        print(f"  ! missing: {p}")
        return None
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    return slide.shapes.add_picture(str(p), x, y, **kw)

def thin_rule(slide, x, y, w, color=RULE, weight=0.75):
    sh = slide.shapes.add_connector(1, x, y, x + w, y)
    sh.line.color.rgb = color
    sh.line.width = Pt(weight)
    return sh

def slide_label(slide, label_text, color=ACCENT):
    text(slide, label_text.upper(),
         Inches(0.6), Inches(0.5), Inches(8), Inches(0.3),
         size=9, bold=True, color=color, font=HEAD)

def slide_number(slide, n, total):
    text(slide, f"{n:02d} / {total:02d}",
         Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
         size=9, color=MUTE, font=MONO, align=PP_ALIGN.RIGHT)


# Deep-link helper — relative paths (relative to PPTX location);
# PowerPoint resolves them when the PPTX is in the same folder as the file.
def rel_link(target: Path) -> str:
    """Make a path relative to the PPTX's location for deep-link hyperlinks."""
    try:
        return str(target.relative_to(OUT_REPO.parent))
    except ValueError:
        # outside the presentation folder — use a relative-up path
        try:
            return str(target.relative_to(REPO).resolve().relative_to(OUT_REPO.parent))
        except Exception:
            return "../../" + str(target.relative_to(REPO))


TOTAL = 25

# ===== 01 Title =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, DARK)
text(slide, "AMINAK   /   PHASE 14",
     Inches(0.7), Inches(0.6), Inches(8), Inches(0.3),
     size=10, bold=True, color=RGBColor(0xC9,0xCC,0xD2), font=HEAD)
text(slide, "Designing inhibitors for TYMS\nat four binding sites.",
     Inches(0.7), Inches(2.0), Inches(11.5), Inches(2.6),
     size=50, bold=True, color=WHITE, font=HEAD)
thin_rule(slide, Inches(0.7), Inches(4.7), Inches(2.0), color=ACCENT, weight=1.5)
text(slide, "Active site · Cofactor site · Dimer interface · Allosteric",
     Inches(0.7), Inches(4.9), Inches(11.5), Inches(0.5),
     size=15, color=RGBColor(0xC9,0xCC,0xD2), font=BODY, italic=True)
# v3 banner
text(slide, "v3 — with Smina rescoring, MM-GBSA + HADDOCK3 roadmap, PPI/BSA, Modeller↔AlphaFold, Plasmodium selectivity",
     Inches(0.7), Inches(5.8), Inches(11.5), Inches(0.5),
     size=11, color=ACCENT, italic=True, font=BODY)
text(slide, [{"text":"github.com/ArioMoniri/aminak", "color":RGBColor(0xC9,0xCC,0xD2)}],
     Inches(0.7), Inches(6.85), Inches(10), Inches(0.3), size=10, font=MONO)
text(slide, "Ariorad Moniri   /   2026",
     Inches(0.7), Inches(7.15), Inches(10), Inches(0.3),
     size=9, color=MUTE, font=HEAD)

# ===== 02 The question (3 stats) =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, BONE)
slide_label(slide, "Setup")
slide_number(slide, 2, TOTAL)
text(slide, "The question.", Inches(0.7), Inches(1.0), Inches(12), Inches(1.0),
     size=44, bold=True, color=INK, font=HEAD)
text(slide, "What molecules will out-compete dUMP at the active site — or bind elsewhere on TYMS?",
     Inches(0.7), Inches(2.3), Inches(12), Inches(0.6),
     size=18, color=CHAR, font=BODY, italic=True)
thin_rule(slide, Inches(0.7), Inches(3.3), Inches(12), color=RULE)
def stat(big, label, x, big_color=INK):
    text(slide, big, x, Inches(3.9), Inches(4.0), Inches(1.6),
         size=82, bold=True, color=big_color, font=HEAD)
    text(slide, label, x, Inches(5.7), Inches(4.0), Inches(0.8),
         size=13, color=CHAR, font=BODY)
stat("4", "binding-site strategies", Inches(0.7))
stat("86", "docked compounds", Inches(4.7))
stat("7", "reviewer / corrector rounds", Inches(8.7), big_color=ACCENT)
text(slide, "TYMS · UniProt P04818 · PDB 1HVY · AutoDock Vina 1.2.7 + Smina · ±0.85 kcal/mol noise floor",
     Inches(0.7), Inches(6.9), Inches(12), Inches(0.3),
     size=10, color=MUTE, font=MONO)

# ===== 03 Pipeline =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Pipeline · 8 phases")
slide_number(slide, 3, TOTAL)
text(slide, "From conservation to inhibitor.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=32, bold=True, color=INK, font=HEAD)
add_image(slide, "workflow", Inches(1.5), Inches(2.0), w=Inches(10.5))
text(slide, "Each phase reused by the next. Each independently reviewed by paired agents.",
     Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
     size=11, color=MUTE, font=BODY, italic=True, align=PP_ALIGN.CENTER)

# ===== 04 Conservation multipanel =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 1+2 · MSA + active site")
slide_number(slide, 4, TOTAL)
text(slide, "The active site is highly conserved.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
add_image(slide, "conservation_multipanel", Inches(0.4), Inches(1.7), w=Inches(12.5))
text(slide, "→ defines the docking box used in every later phase.",
     Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
     size=12, italic=True, color=ACCENT, font=HEAD, bold=True, align=PP_ALIGN.CENTER)

# ===== 05 Structure prep — cleaner images =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 3+4 · Structure")
slide_number(slide, 5, TOTAL)
text(slide, "Dimer-aware receptor preparation.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
add_image(slide, "dimer_overview_clean", Inches(0.5), Inches(1.9), w=Inches(4.2))
add_image(slide, "dimer_activesite_clean", Inches(4.8), Inches(1.9), w=Inches(4.2))
add_image(slide, "cavity18_carve_clean",   Inches(9.0), Inches(1.9), w=Inches(4.0))
text(slide, "Dimer overview", Inches(0.5), Inches(5.5), Inches(4.2), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Active-site close-up", Inches(4.8), Inches(5.5), Inches(4.2), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Cavity 18 carve-out (wheat surface)", Inches(9.0), Inches(5.5), Inches(4.0), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Both chains kept · AMBER ff14SB charges (pdb2pqr30 + PROPKA pH 7.4) · cofactor net −2 e per copy",
     Inches(0.7), Inches(6.05), Inches(12), Inches(0.5),
     size=12, color=CHAR, font=BODY, align=PP_ALIGN.CENTER)
# Deep-link
text(slide,
     [{"text": "▸ download the prepared dimer PDB · ",
       "color": MUTE, "size": 10, "italic": True},
      {"text": "06f_receptor_fixed/protein_dimer_apo_fixed.pdbqt",
       "color": LINK, "size": 10, "font": MONO,
       "hyperlink": "../../06f_receptor_fixed/protein_dimer_apo_fixed.pdbqt"}],
     Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
     align=PP_ALIGN.CENTER)

# ===== 06 WT docking — CORRECTED =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 5+6 · WT Vina docking")
slide_number(slide, 6, TOTAL)
text(slide, "Wild-type dUMP docks at −8.55 ± 0.05 kcal/mol.",
     Inches(0.7), Inches(0.95), Inches(12.5), Inches(0.7),
     size=26, bold=True, color=INK, font=HEAD)
text(slide, "5-seed mean (Phase 7 multi-replica protocol) · single-seed 42 = −8.78 (Phase 6 canonical)",
     Inches(0.7), Inches(1.6), Inches(12), Inches(0.4),
     size=12, color=MUTE, italic=True, font=BODY)

text(slide, "−8.55", Inches(0.7), Inches(2.3), Inches(6.0), Inches(2.0),
     size=110, bold=True, color=INK, font=HEAD)
text(slide, "kcal/mol  (mean of 5 seeds)", Inches(0.7), Inches(4.6), Inches(6.0), Inches(0.4),
     size=15, color=MUTE, italic=True, font=BODY)

text(slide, ["Phase 7 multi-replica spread", "−8.59 to −8.47 across {42, 7, 13, 99, 256}",
             "SD = 0.05 kcal/mol  (sub-noise variance)",
             "Holo state (cofactor present): −7.49 kcal/mol"],
     Inches(0.7), Inches(5.3), Inches(6.0), Inches(2.0),
     size=13, color=CHAR, font=BODY)

add_image(slide, "holo_dump_cofactor", Inches(7.5), Inches(1.9), w=Inches(5.3))
text(slide, "Holo state (dUMP + raltitrexed) →", Inches(7.5), Inches(6.4), Inches(5.3), Inches(0.4),
     size=12, italic=True, color=MUTE, font=BODY)

# ===== 07 Phase 7 null =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 7 · mutagenesis", color=ACCENT)
slide_number(slide, 7, TOTAL)
text(slide, "Rigid Vina cannot resolve mutants.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
add_image(slide, "delta_apo", Inches(0.7), Inches(2.0), w=Inches(7.5))
text(slide, "+0.77", Inches(8.7), Inches(2.4), Inches(4.0), Inches(1.4),
     size=80, bold=True, color=ACCENT, font=HEAD)
text(slide, "largest holo Δ Vina (kcal/mol)",
     Inches(8.7), Inches(4.0), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "±0.85", Inches(8.7), Inches(4.7), Inches(4.0), Inches(1.0),
     size=60, bold=True, color=INK, font=HEAD)
text(slide, "Vina noise floor (Trott & Olson 2010)",
     Inches(8.7), Inches(6.0), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "→ honest null result.",
     Inches(8.7), Inches(6.6), Inches(4.0), Inches(0.4),
     size=14, italic=True, bold=True, color=ACCENT, font=HEAD)

# ===== 08 Phase 8 smarter scoring =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 8 · Vinardo + flex-Vina")
slide_number(slide, 8, TOTAL)
text(slide, "Can smarter scoring fix it? — Partially.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=28, bold=True, color=INK, font=HEAD)
add_image(slide, "phase8_flex", Inches(0.7), Inches(2.0), w=Inches(6.0))
add_image(slide, "phase8_alt",  Inches(7.0), Inches(2.0), w=Inches(5.8))
text(slide, "Flex-residue Vina", Inches(0.7), Inches(6.0), Inches(6.0), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Vinardo vs Vina", Inches(7.0), Inches(6.0), Inches(5.8), Inches(0.3),
     size=11, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "Partial fix only. The R→E sign error needs proper PB electrostatics.",
     Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
     size=13, color=CHAR, font=BODY, italic=True, align=PP_ALIGN.CENTER)

# ===== 09 Modeller homology — Ramachandran improvement =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 6 · Modeller + Ramachandran")
slide_number(slide, 9, TOTAL)
text(slide, "Reducing Ramachandran outliers, step by step.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=26, bold=True, color=INK, font=HEAD)
text(slide, "Phase 6 → 6b: validator + refinement upgrades",
     Inches(0.7), Inches(1.65), Inches(12), Inches(0.4),
     size=13, color=MUTE, italic=True, font=BODY)

# Left: process bullets
text(slide, "What we did", Inches(0.7), Inches(2.1), Inches(5.0), Inches(0.4),
     size=14, bold=True, color=DEEP if "DEEP" in dir() else INK, font=HEAD)
text(slide, [
    [{"text": "1. Lovell 4-map validator ", "bold": True, "color": INK},
     {"text": "(general / Gly / Pro / pre-Pro)  — replaced the single-polygon validator;", "color": CHAR}],
    [{"text": "2. md_level = refine.very_slow ", "bold": True, "color": INK},
     {"text": "+ max_var_iterations = 600 + repeat_optimization = 2;", "color": CHAR}],
    [{"text": "3. LoopModel on residues 93–101 ", "bold": True, "color": INK},
     {"text": "(template-disagreement region);", "color": CHAR}],
    [{"text": "4. Ensemble of 10 models", "bold": True, "color": INK},
     {"text": " — best chosen by Lovell %favoured.", "color": CHAR}],
], Inches(0.7), Inches(2.55), Inches(5.4), Inches(3.5), size=11)

# Right: image showing before/after
add_image(slide, "rama_compare", Inches(6.5), Inches(2.0), w=Inches(6.5))

# Bottom: the journey of %favoured
text(slide, "Best-model % favoured (Lovell)",
     Inches(0.7), Inches(6.1), Inches(12), Inches(0.4),
     size=13, bold=True, color=INK, font=HEAD)
text(slide, [
    [{"text":"baseline", "color":MUTE}, {"text":"  →  ", "color":MUTE},
     {"text":"83.5–85.3 %", "bold":True}, {"text":"  (broken validator)", "color":MUTE, "italic":True},
     {"text":"  →  ", "color":MUTE},
     {"text":"+ Lovell partition", "color":MUTE}, {"text":"  →  ", "color":MUTE},
     {"text":"94.7–96.1 %", "bold":True, "color":GREEN},
     {"text":"  →  ", "color":MUTE},
     {"text":"+ refine.very_slow", "color":MUTE}, {"text":"  →  ", "color":MUTE},
     {"text":"95.4 %", "bold":True, "color":GREEN},
     {"text":"   (1HVY crystal: 92.2 %)", "color":MUTE, "italic":True}],
], Inches(0.7), Inches(6.5), Inches(12), Inches(0.5), size=11)

# ===== 10 Modeller vs AlphaFold =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 7c · Modeller ↔ AlphaFold")
slide_number(slide, 10, TOTAL)
text(slide, "Modeller vs AlphaFold — do the two methods agree?",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=24, bold=True, color=INK, font=HEAD)
add_image(slide, "modeller_af", Inches(0.7), Inches(1.7), w=Inches(12.0))
text(slide, "Refined Modeller (#3, #10) and AlphaFold v6 all sit within 0.4 Å Cα RMSD of the 1HVY crystal — and all 3 beat the crystal on Lovell %favoured (94.5–95.4 % vs 92.2 %).",
     Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
     size=12, color=CHAR, font=BODY, italic=True, align=PP_ALIGN.CENTER)

# ===== 11 PPI / BSA dimer interface =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 14 · PPI / BSA", color=ACCENT)
slide_number(slide, 11, TOTAL)
text(slide, "Dimer interface — the same residues bind dUMP and stabilise the dimer.",
     Inches(0.7), Inches(0.95), Inches(12.5), Inches(0.7),
     size=22, bold=True, color=INK, font=HEAD)
add_image(slide, "ppi_dimer", Inches(0.4), Inches(1.7), w=Inches(9.0))
# Right column
text(slide, "Total BSA", Inches(9.6), Inches(1.9), Inches(3.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, "≈ 4 160", Inches(9.6), Inches(2.3), Inches(3.5), Inches(1.2),
     size=58, bold=True, color=INK, font=HEAD)
text(slide, "Å² total (2 079 per side)",
     Inches(9.6), Inches(3.6), Inches(3.5), Inches(0.4),
     size=11, color=CHAR, font=BODY)
text(slide, "Top hot-spots — chain A",
     Inches(9.6), Inches(4.2), Inches(3.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, [
    [{"text":"R175  ", "bold":True, "color":ACCENT}, {"text":"184.5 Å²"}],
    [{"text":"P59   ", "bold":True}, {"text":"156.7 Å²"}],
    [{"text":"R176  ", "bold":True, "color":ACCENT}, {"text":"137.9 Å²"}],
    [{"text":"R202  ", "bold":True}, {"text":"116.3 Å²"}],
    [{"text":"R215  ", "bold":True, "color":ACCENT}, {"text":" 66.9 Å²"}],
], Inches(9.6), Inches(4.6), Inches(3.5), Inches(2.0), size=11, color=CHAR, font=BODY)
text(slide, "★ R175 / R176 / R215 are also the phosphate-clamp residues — substrate + dimer share the same hot-spots.",
     Inches(0.7), Inches(6.9), Inches(12.5), Inches(0.4),
     size=11, italic=True, color=ACCENT, font=BODY, align=PP_ALIGN.CENTER)

# ===== 12 Section divider — Phase 14 =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, DARK)
text(slide, "PHASE 14", Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
     size=12, bold=True, color=ACCENT, font=HEAD)
text(slide, "Inhibitor design.",
     Inches(0.7), Inches(2.8), Inches(12), Inches(1.6),
     size=64, bold=True, color=WHITE, font=HEAD)
thin_rule(slide, Inches(0.7), Inches(4.5), Inches(3.0), color=ACCENT, weight=1.5)
text(slide, "Four orthogonal sites. 86 docked compounds. Two findings.",
     Inches(0.7), Inches(4.8), Inches(12), Inches(0.5),
     size=18, color=RGBColor(0xC9,0xCC,0xD2), italic=True, font=BODY)

# ===== 13 S1 active site =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 1 · active site")
slide_number(slide, 13, TOTAL)
text(slide, "Clean active-vs-prodrug gap. Kcal-noise silent.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=26, bold=True, color=INK, font=HEAD)
add_image(slide, "phase14_tier", Inches(0.7), Inches(2.0), w=Inches(8.0))
text(slide, "−9.04", Inches(9.0), Inches(2.3), Inches(4.0), Inches(1.4),
     size=72, bold=True, color=INK, font=HEAD)
text(slide, "5-FdUMP   (canonical)",
     Inches(9.0), Inches(3.8), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "−4.95", Inches(9.0), Inches(4.6), Inches(4.0), Inches(1.4),
     size=72, bold=True, color=MUTE, font=HEAD)
text(slide, "5-FU   (prodrug)",
     Inches(9.0), Inches(6.1), Inches(4.0), Inches(0.4),
     size=12, color=CHAR, font=BODY)
text(slide, "≈3 kcal/mol enrichment signal.",
     Inches(9.0), Inches(6.6), Inches(4.0), Inches(0.4),
     size=11, italic=True, bold=True, color=ACCENT, font=HEAD)

# ===== 14 S2 Plevitrexed =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 2 · cofactor site", color=ACCENT)
slide_number(slide, 14, TOTAL)
text(slide, "Plevitrexed (ZD9331) above noise.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
add_image(slide, "phase14_delta", Inches(0.7), Inches(2.0), w=Inches(7.5))
text(slide, "−10.01", Inches(8.6), Inches(2.4), Inches(4.5), Inches(1.6),
     size=92, bold=True, color=ACCENT, font=HEAD)
text(slide, "kcal/mol", Inches(8.6), Inches(4.2), Inches(4.5), Inches(0.4),
     size=14, color=MUTE, italic=True, font=BODY)
thin_rule(slide, Inches(8.6), Inches(4.8), Inches(2.5), color=INK, weight=1.5)
text(slide, ["Δ −0.88 vs raltitrexed",
             "above the ±0.85 noise floor",
             "reproducible across both seeds",
             "consistent with Ki ~nM (Jackman 1997)"],
     Inches(8.6), Inches(5.0), Inches(4.5), Inches(2.0),
     size=13, color=CHAR, font=BODY)

# ===== 15 S3 dimer null + HADDOCK3 roadmap inline =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 3 · dimer interface")
slide_number(slide, 15, TOTAL)
text(slide, "Documented null → HADDOCK3 is the right next step.",
     Inches(0.7), Inches(0.95), Inches(12.5), Inches(0.7),
     size=26, bold=True, color=INK, font=HEAD)
text(slide, "Rigid Vina cannot resolve ≥6-mer peptides (Hassan 2017). The Phase-14 fragment-decomposition baseline:",
     Inches(0.7), Inches(1.6), Inches(12), Inches(0.4),
     size=12, color=MUTE, italic=True, font=BODY)
text(slide, "Canonical LR-octapeptide   worse   than scrambled control by 1.48 kcal/mol",
     Inches(0.7), Inches(2.1), Inches(12.5), Inches(0.6),
     size=18, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
thin_rule(slide, Inches(0.7), Inches(2.8), Inches(12), color=RULE)

# Bottom half: HADDOCK3 roadmap
text(slide, "HADDOCK3 — the right tool (set up, ready to run)",
     Inches(0.7), Inches(3.1), Inches(12), Inches(0.5),
     size=16, bold=True, color=ACCENT, font=HEAD)
text(slide, [
    [{"text":"▸ Peptide ≥ 6 residues: ", "bold": True}, {"text":"flexible refinement via CNS + ambiguous interaction restraints (AIRs)."}],
    [{"text":"▸ Active residues from Phase-14 A3 contact map: ", "bold": True}, {"text":"34 chain-A residues at the dimer interface."}],
    [{"text":"▸ Passive residues: ", "bold": True}, {"text":"chain-A residues within 6.5 Å of the active set."}],
    [{"text":"▸ Pipeline modules: ", "bold": True}, {"text":"topoaa → rigidbody (1000 decoys) → flexref (200) → mdref (50) → caprieval."}],
    [{"text":"▸ Scrambled control: ", "bold": True}, {"text":"same protocol, numpy seed 42 permutation; specificity = top-cluster Δ."}],
    [{"text":"▸ Wall-time estimate: ", "bold": True}, {"text":"~6 h per peptide on a single Mac (CNS install required)."}],
], Inches(0.7), Inches(3.6), Inches(12.5), Inches(3.0), size=12)
# deep-link
text(slide,
     [{"text":"▸ ready-to-run config & restraints  ·  ", "color":MUTE, "italic":True, "size":10},
      {"text":"14_inhibitor_design/07_advanced_methods/haddock3/", "color":LINK, "size":10, "font":MONO,
       "hyperlink":"../07_advanced_methods/haddock3/README.md"}],
     Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))

# ===== 16 S4 cavity 18 headline =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Strategy 4 · allosteric", color=ACCENT)
slide_number(slide, 16, TOTAL)
text(slide, "An under-explored druggable cavity.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=30, bold=True, color=INK, font=HEAD)
def big_stat(x, big, label, color=INK):
    text(slide, big, x, Inches(2.4), Inches(3.0), Inches(1.6),
         size=70, bold=True, color=color, font=HEAD)
    text(slide, label, x, Inches(4.0), Inches(3.0), Inches(0.8),
         size=11, color=CHAR, font=BODY)
big_stat(Inches(0.7),  "0.994", "FPocket druggability\ncavity 18", color=ACCENT)
big_stat(Inches(3.8),  "0.828", "C2-symmetric mirror\ncavity 17", color=INK)
big_stat(Inches(6.9),  "−7.52", "1H-indazole\nkcal/mol", color=ACCENT)
big_stat(Inches(10.0), "−7.28", "ibuprofen\nkcal/mol", color=ACCENT)
thin_rule(slide, Inches(0.7), Inches(5.6), Inches(12), color=RULE)
text(slide, "Two unrelated drug-like fragments. Same pocket. Different chemistries.",
     Inches(0.7), Inches(5.9), Inches(12), Inches(0.5),
     size=18, bold=True, color=INK, font=HEAD, align=PP_ALIGN.CENTER)
text(slide, "≈ 2 kcal/mol better than anywhere else on the chain-A or chain-B surface.",
     Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
     size=13, color=MUTE, font=BODY, italic=True, align=PP_ALIGN.CENTER)

# ===== 17 indazole pose =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · top hit", color=ACCENT)
slide_number(slide, 17, TOTAL)
text(slide, "1H-indazole.",
     Inches(0.7), Inches(0.95), Inches(8), Inches(0.7),
     size=36, bold=True, color=INK, font=HEAD)
text(slide, "Kinase-inhibitor scaffold · −7.52 kcal/mol",
     Inches(0.7), Inches(1.7), Inches(8), Inches(0.5),
     size=14, color=MUTE, italic=True, font=BODY)
add_image(slide, "pose_indazole", Inches(0.5), Inches(2.4), w=Inches(7.5))
text(slide, "Engages", Inches(8.5), Inches(2.4), Inches(4.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, [
    [{"text": "Phe55  ", "bold": True, "color": INK}, {"text": "H-bond + π"}],
    [{"text": "Asn201  ", "bold": True, "color": INK}, {"text": "H-bond"}],
    [{"text": "Phe200  ", "bold": True, "color": INK}, {"text": "π-stack"}],
    [{"text": "Leu196   ", "bold": True, "color": ACCENT}, {"text": "allosteric loop 181-197 ★", "color": ACCENT}],
    [{"text": "Gly197   ", "bold": True, "color": ACCENT}, {"text": "allosteric loop 181-197 ★", "color": ACCENT}],
    [{"text": "Ile83  ", "bold": True, "color": INK}, {"text": "hydrophobic wall"}],
], Inches(8.5), Inches(2.85), Inches(4.5), Inches(3.5), size=14, color=CHAR, font=BODY)
text(slide,
     [{"text":"▸ open 3D viewer  ·  ", "color":MUTE, "italic":True, "size":10},
      {"text":"cavity18_indazole.html", "color":LINK, "size":10, "font":MONO,
       "hyperlink":"../04_allosteric/cavity18_evidence/viewers/cavity18_indazole.html"}],
     Inches(8.5), Inches(6.3), Inches(4.5), Inches(0.4))
text(slide, "3 of 10 contact residues are on the published allosteric communication loop (Anderson 2012, Pozzi 2019).",
     Inches(8.5), Inches(6.7), Inches(4.5), Inches(0.7),
     size=10, italic=True, color=ACCENT, font=BODY)

# ===== 18 ibuprofen pose =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · second hit", color=ACCENT)
slide_number(slide, 18, TOTAL)
text(slide, "Ibuprofen.",
     Inches(0.7), Inches(0.95), Inches(8), Inches(0.7),
     size=36, bold=True, color=INK, font=HEAD)
text(slide, "NSAID, COX1/2 · −7.28 kcal/mol",
     Inches(0.7), Inches(1.7), Inches(8), Inches(0.5),
     size=14, color=MUTE, italic=True, font=BODY)
add_image(slide, "pose_ibuprofen", Inches(0.5), Inches(2.4), w=Inches(7.5))
text(slide, "Engages", Inches(8.5), Inches(2.4), Inches(4.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, [
    [{"text": "Lys283  ", "bold": True, "color": ACCENT}, {"text": "salt bridge ★", "color": ACCENT}],
    [{"text": "Lys52   ", "bold": True, "color": ACCENT}, {"text": "salt bridge ★", "color": ACCENT}],
    [{"text": "Phe200  ", "bold": True, "color": INK}, {"text": "π-stack"}],
    [{"text": "Leu196  ", "bold": True, "color": ACCENT}, {"text": "loop 181-197 ★", "color": ACCENT}],
    [{"text": "Phe55   ", "bold": True, "color": INK}, {"text": "hydrophobic"}],
    [{"text": "Ile83   ", "bold": True, "color": INK}, {"text": "hydrophobic"}],
], Inches(8.5), Inches(2.85), Inches(4.5), Inches(3.5), size=14, color=CHAR, font=BODY)
text(slide,
     [{"text":"▸ open 3D viewer  ·  ", "color":MUTE, "italic":True, "size":10},
      {"text":"cavity18_ibuprofen.html", "color":LINK, "size":10, "font":MONO,
       "hyperlink":"../04_allosteric/cavity18_evidence/viewers/cavity18_ibuprofen.html"}],
     Inches(8.5), Inches(6.3), Inches(4.5), Inches(0.4))
text(slide, "Double-lysine salt-bridge clamps the deprotonated carboxylate at pH 7.4.",
     Inches(8.5), Inches(6.7), Inches(4.5), Inches(0.7),
     size=10, italic=True, color=ACCENT, font=BODY)

# ===== 19 Conservation cavity 18 =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · conservation", color=ACCENT)
slide_number(slide, 19, TOTAL)
text(slide, "Conserved at the contact face.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=28, bold=True, color=INK, font=HEAD)
add_image(slide, "conservation_cav18", Inches(0.7), Inches(2.0), w=Inches(8.5))
text(slide, "100% conserved", Inches(9.5), Inches(2.3), Inches(3.5), Inches(0.4),
     size=11, bold=True, color=MUTE, font=HEAD)
text(slide, "across 10 orthologs",
     Inches(9.5), Inches(2.6), Inches(3.5), Inches(0.3),
     size=10, italic=True, color=MUTE, font=BODY)
text(slide, [
    [{"text": "Gly54", "bold": True, "color": INK}],
    [{"text": "Glu87", "bold": True, "color": INK}],
    [{"text": "Met190", "bold": True, "color": ACCENT}, {"text": "  loop", "color": MUTE, "size": 10}],
    [{"text": "Ala191", "bold": True, "color": ACCENT}, {"text": "  loop", "color": MUTE, "size": 10}],
    [{"text": "Leu196", "bold": True, "color": ACCENT}, {"text": "  loop", "color": MUTE, "size": 10}],
    [{"text": "Phe200", "bold": True, "color": INK}],
    [{"text": "Asn201", "bold": True, "color": INK}],
], Inches(9.5), Inches(3.1), Inches(3.5), Inches(2.8), size=14, color=CHAR, font=BODY)
text(slide, "6 of 7 are contact residues for at least one top hit.",
     Inches(9.5), Inches(6.2), Inches(3.5), Inches(1.0),
     size=11, italic=True, bold=True, color=ACCENT, font=BODY)

# ===== 20 Phylogeny + selectivity explained =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Cavity 18 · selectivity", color=ACCENT)
slide_number(slide, 20, TOTAL)
text(slide, "Why a Plasmodium-selective drug is plausible here.",
     Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
     size=24, bold=True, color=INK, font=HEAD)
add_image(slide, "phylogeny_cav18", Inches(0.5), Inches(1.7), w=Inches(7.5))

# Explanation
text(slide, "Plain-English meaning", Inches(8.3), Inches(1.9), Inches(4.7), Inches(0.4),
     size=13, bold=True, color=ACCENT, font=HEAD)
text(slide, [
    [{"text":"Malaria parasites use TYMS too. ", "color":INK, "bold":True},
     {"text":"To kill the parasite without poisoning the patient, a drug must hit ", "color":CHAR},
     {"text":"Plasmodium ", "italic":True, "color":CHAR},
     {"text":"TYMS but spare human TYMS.", "color":CHAR}],
    [{"text":"The classical active site is too conserved ", "color":INK, "bold":True},
     {"text":"— 5-FU / raltitrexed hit both equally → toxicity.", "color":CHAR}],
    [{"text":"Cavity 18 differs at 21 of 36 residues between human and ", "color":INK, "bold":True},
     {"text":"P. falciparum", "italic":True, "bold":True, "color":INK},
     {"text":". A cavity-18 ligand could plausibly be tuned to grip the parasite shape while sliding past the human one.", "color":CHAR}],
    [{"text":"This is the kind of species-selectivity handle ", "color":INK, "bold":True},
     {"text":"that anti-parasitic medicinal chemistry actively searches for. ", "color":CHAR},
     {"text":"Hypothesis from the mutation distribution — pending experimental confirmation.", "italic":True, "color":MUTE}],
], Inches(8.3), Inches(2.35), Inches(4.7), Inches(4.5), size=11)

# ===== 21 Smina rescoring (electrostatics) — the negative finding =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 14e · Smina rescoring")
slide_number(slide, 21, TOTAL)
text(slide, "Adding electrostatics ≠ fixing the R→E sign error.",
     Inches(0.7), Inches(0.95), Inches(12.5), Inches(0.7),
     size=24, bold=True, color=INK, font=HEAD)
text(slide, "Smina (Vina + Coulomb + AD4 desolvation + custom weights) tested on the full Phase-7 panel.",
     Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
     size=12, color=MUTE, italic=True, font=BODY)
add_image(slide, "smina_rescore", Inches(0.5), Inches(2.0), w=Inches(12.3))
text(slide,
     [{"text":"▸ details + custom scoring files  ·  ", "color":MUTE, "italic":True, "size":10},
      {"text":"06_smina_rescore/", "color":LINK, "size":10, "font":MONO,
       "hyperlink":"../06_smina_rescore/"},
      {"text":"   ★ Smina DOES capture cavity-18 ibuprofen salt-bridge (q_amp Δ = −4.4 vs indazole)",
       "color":ACCENT, "italic":True, "size":11}],
     Inches(0.7), Inches(6.8), Inches(12.5), Inches(0.4))

# ===== 22 OpenMM MM-GBSA equivalent — ACTUALLY EXECUTED =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Phase 14g · OpenMM GB rescoring — ★ executed", color=ACCENT)
slide_number(slide, 22, TOTAL)
text(slide, "Side-chain relaxation + GB electrostatics exposes the sign error.",
     Inches(0.7), Inches(0.95), Inches(12.5), Inches(0.7),
     size=22, bold=True, color=INK, font=HEAD)
text(slide, "AMBER ff14SB + GBn2 implicit · 5 mutants minimised on arm64 CPU · ~2 min/system",
     Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
     size=12, color=MUTE, italic=True, font=BODY)
add_image(slide, "openmm_gb", Inches(0.5), Inches(2.0), w=Inches(7.5))
# Right column: stats
text(slide, "Rank order respects physics", Inches(8.5), Inches(2.0), Inches(4.5), Inches(0.4),
     size=14, bold=True, color=ACCENT, font=HEAD)
text(slide, [
    [{"text": "R175E_R176E ", "bold": True, "color": ACCENT}, {"text": "  +328  ★ max"}],
    [{"text": "R215A       ", "bold": True}, {"text": "  +165"}],
    [{"text": "R215E       ", "bold": True}, {"text": "  +158"}],
    [{"text": "R175E       ", "bold": True}, {"text": "  +132"}],
    [{"text": "C195A       ", "bold": True}, {"text": "  +61   (no charge change)"}],
], Inches(8.5), Inches(2.5), Inches(4.5), Inches(2.5), size=12, color=CHAR, font=MONO)
text(slide, "Δ E_receptor vs WT_holo (kcal/mol).",
     Inches(8.5), Inches(5.1), Inches(4.5), Inches(0.4),
     size=11, italic=True, color=MUTE, font=BODY)
text(slide, "DOUBLE charge reversal = largest penalty. Rigid Vina / Smina could not see this.",
     Inches(8.5), Inches(5.5), Inches(4.5), Inches(1.5),
     size=12, bold=True, italic=True, color=ACCENT, font=BODY)
text(slide,
     [{"text":"▸ source data  ·  ", "color":MUTE, "italic":True, "size":10},
      {"text":"openmm_gb_results.csv", "color":LINK, "size":10, "font":MONO,
       "hyperlink":"../07_advanced_methods/openmm_gb_rescore/openmm_gb_results.csv"}],
     Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))
text(slide, "What the OpenMM-GB executed run already showed:", Inches(0.7), Inches(2.1), Inches(12), Inches(0.4),
     size=14, bold=True, color=ACCENT, font=HEAD)
text(slide, [
    "double Arg→Glu reversal R175E_R176E is the largest penalty (+328 kcal/mol)",
    "the rank — double > singles > neutralisation > catalytic — is physically correct",
    "rigid Vina / Smina (Phase 14e) could NOT produce this signal",
    "the missing physics IS structural relaxation + implicit-solvent electrostatics",
], Inches(0.7), Inches(2.55), Inches(12.5), Inches(2.0), size=12, color=CHAR, font=BODY)

text(slide, "What still needs the full AmberTools MMPBSA.py path:", Inches(0.7), Inches(4.6), Inches(12), Inches(0.4),
     size=14, bold=True, color=INK, font=HEAD)
text(slide, [
    "absolute ΔΔG_bind (not just rank order) — needs ligand GAFF parametrisation + thermodynamic cycle",
    "MD ensemble averaging — single-pose enthalpy includes relaxation strain",
    "Poisson-Boltzmann (not just GB) electrostatics for high-salt regimes",
    "compute estimate: 24 h/mutant on CPU sander or overnight for the sweep on a CUDA box",
], Inches(0.7), Inches(5.0), Inches(12.5), Inches(2.0), size=12, color=CHAR, font=BODY)

text(slide, "Compute estimate", Inches(0.7), Inches(5.4), Inches(12), Inches(0.4),
     size=13, bold=True, color=INK, font=HEAD)
text(slide, [
    [{"text":"Without GPU: ", "bold":True}, {"text":"~24 h per mutant on CPU sander → ~10 days for the 10-system sweep."}],
    [{"text":"With CUDA box: ", "bold":True}, {"text":"overnight for the whole sweep with pmemd.cuda."}],
], Inches(0.7), Inches(5.85), Inches(12.5), Inches(1.0), size=11)
text(slide,
     [{"text":"▸ ready-to-run scripts  ·  ", "color":MUTE, "italic":True, "size":10},
      {"text":"07_advanced_methods/mmgbsa/", "color":LINK, "size":10, "font":MONO,
       "hyperlink":"../07_advanced_methods/mmgbsa/README.md"}],
     Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))

# ===== 23 Full Smina panel (showing the negative finding) =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, PAPER)
slide_label(slide, "Smina · full Phase-7 panel")
slide_number(slide, 23, TOTAL)
text(slide, "Same negative finding on every mutant.",
     Inches(0.7), Inches(0.95), Inches(12.5), Inches(0.7),
     size=24, bold=True, color=INK, font=HEAD)
text(slide, "42 Phase-7 mutant top poses × 4 scorers (Vina · Vinardo · custom_q · q_amp).",
     Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
     size=12, color=MUTE, italic=True, font=BODY)
add_image(slide, "smina_full", Inches(0.5), Inches(2.05), w=Inches(12.3))

# ===== 24 TL;DR =====
slide = prs.slides.add_slide(BLANK); add_bg(slide, DARK)
text(slide, "TL;DR", Inches(0.7), Inches(0.6), Inches(8), Inches(0.4),
     size=12, bold=True, color=ACCENT, font=HEAD)
text(slide, "We didn't find a new clinical lead.",
     Inches(0.7), Inches(1.5), Inches(12), Inches(0.9),
     size=38, bold=True, color=WHITE, font=HEAD)
text(slide, "We found a plausible third TYMS binding site.",
     Inches(0.7), Inches(2.5), Inches(12), Inches(0.9),
     size=38, bold=True, color=ACCENT, font=HEAD)
thin_rule(slide, Inches(0.7), Inches(3.7), Inches(2.5), color=ACCENT, weight=1.5)
text(slide, [
    "FPocket druggability 0.994 + 0.828 C2 mirror.",
    "Indazole + ibuprofen dock at −7.5 kcal/mol.",
    "6 of 7 ultra-conserved residues are at the contact face.",
    "21-position divergence in Plasmodium → species-selective handle.",
    "Smina rescoring confirms the salt-bridge chemistry — and rules out a cheap fix for the R→E sign error.",
    "MM-GBSA + HADDOCK3 are signposted with ready-to-run configs.",
], Inches(0.7), Inches(4.0), Inches(12), Inches(3.0), size=17, color=RGBColor(0xC9,0xCC,0xD2), font=BODY)
text(slide, "github.com/ArioMoniri/aminak  ·  v3 deck",
     Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
     size=10, color=MUTE, font=MONO)

# Save
OUT_REPO.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_REPO)
shutil.copyfile(OUT_REPO, OUT_USER)
print(f"→ {OUT_REPO}  ({OUT_REPO.stat().st_size/1024:.0f} KB)")
print(f"→ {OUT_USER}")
